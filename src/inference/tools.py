#!/usr/bin/env python3
"""
Tool Execution Layer for GAIA Mini Dataset
Implements the executable tool layer used by the GAIA pipeline.

Active Tools (Core):
1. web_search - DuckDuckGo search
2. excel_reader - Pandas-based Excel/CSV reader
3. audio_transcription - OpenAI Whisper (lazy loaded)
4. file_reader - DOCX, JSON, JSONLD, XML, TXT reader
5. python_executor - Safe code execution with pandas/numpy
6. pptx_reader - PowerPoint reader with slide markers
7. pdf_reader - PDF text extraction
8. zip_extractor - ZIP archive extraction
9. download_file - URL file downloader
10. web_browser - requests + BeautifulSoup
11. calculator - Python expression evaluator
12. submit_final_answer - Submit the final answer
13. image_recognition - Vision-model-backed image analysis

Placeholder / no-op tools:
video_analysis, code_interpreter

"""

import os
import sys
import json
import io
import difflib
import mimetypes
import re
import signal
import subprocess
import zipfile
import tempfile
import shutil
from pathlib import Path
from typing import Dict, Any, Optional, List, Tuple
from urllib.parse import urlparse, unquote
from src.config.models import (
    get_api_base,
    get_api_key,
    normalize_api_base,
    supports_native_vision_input,
)


def _chat_completion_token_kwargs(model: str, n: int) -> Dict[str, int]:
    """Return the right token-budget kwarg for chat.completions.

    OpenAI's gpt-5 family and o1/o3/o4 reasoning models reject `max_tokens`
    and require `max_completion_tokens`. Everything else (including
    Some OpenAI-compatible gateways keep the legacy `max_tokens`.
    """
    if model:
        m = model.lower()
        if m.startswith(("gpt-5", "o1", "o3", "o4")):
            return {"max_completion_tokens": n}
    return {"max_tokens": n}


def _chat_completion_sampling_kwargs(model: str) -> Dict[str, float]:
    """Return sampling controls accepted by chat.completions for this model."""
    m = (model or "").strip().lower()
    if m == "gpt-5" or m.startswith(("gpt-5-", "gpt-5.5", "o1", "o3", "o4")):
        return {}
    return {"temperature": 0.0}


# Lazy-loaded modules
_whisper_model = None


_DEFAULT_HTTP_HEADERS = {
    "User-Agent": (
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
        "AppleWebKit/537.36 (KHTML, like Gecko) "
        "Chrome/122.0.0.0 Safari/537.36"
    ),
    "Accept-Language": "en-US,en;q=0.9",
    "Accept": "*/*",
}

DEFAULT_ARTIFACT_DIR = os.path.join("runtime_artifacts", "extracted_data")

_MAX_ATTACHMENT_CANDIDATES = 5000


def _find_tesseract_binary() -> str:
    """Find the tesseract executable even when the env is invoked by absolute python."""
    for candidate in (
        shutil.which("tesseract"),
        os.path.join(os.environ.get("CONDA_PREFIX", ""), "bin", "tesseract"),
        os.path.join(os.path.dirname(sys.executable), "tesseract"),
    ):
        if candidate and os.path.exists(candidate):
            return candidate
    return "tesseract"


def _is_http_url(value: str) -> bool:
    """Return True when a string looks like an HTTP(S) URL."""
    if not value:
        return False
    parsed = urlparse(str(value))
    return parsed.scheme in {"http", "https"} and bool(parsed.netloc)


def _normalize_download_url(url: str) -> str:
    """
    Normalize common downloadable URLs without changing the external tool API.

    This helps existing tools cooperate better on sources such as arXiv, where
    a model may give an abstract page URL but downstream readers work better
    with the PDF endpoint.
    """
    if not _is_http_url(url):
        return url

    parsed = urlparse(url)
    host = parsed.netloc.lower()
    path = parsed.path

    if "arxiv.org" in host and path.startswith("/abs/"):
        paper_id = path.split("/abs/", 1)[1].strip("/")
        if paper_id:
            return f"{parsed.scheme}://{parsed.netloc}/pdf/{paper_id}.pdf"

    return url


def _looks_like_pdf_url(url: str) -> bool:
    """Heuristic check for URLs that should be treated as PDFs."""
    if not _is_http_url(url):
        return False
    normalized = _normalize_download_url(url)
    path = urlparse(normalized).path.lower()
    return path.endswith(".pdf") or "/pdf/" in path


def _http_headers(accept: Optional[str] = None, referer: Optional[str] = None) -> Dict[str, str]:
    """Create consistent browser-like HTTP headers."""
    headers = dict(_DEFAULT_HTTP_HEADERS)
    if accept:
        headers["Accept"] = accept
    if referer:
        headers["Referer"] = referer
    return headers


def _guess_filename_from_response(url: str, response) -> str:
    """Infer a filename from response headers or URL."""
    content_disposition = response.headers.get("Content-Disposition", "")
    if "filename=" in content_disposition:
        filename = content_disposition.split("filename=", 1)[1].strip().strip('"')
        filename = filename.split(";", 1)[0].strip()
        if filename:
            return filename

    final_url = response.url or url
    filename = unquote(os.path.basename(urlparse(final_url).path))
    if not filename:
        filename = "downloaded_file"

    content_type = response.headers.get("Content-Type", "").split(";", 1)[0].strip().lower()
    if not Path(filename).suffix and content_type:
        guessed_ext = mimetypes.guess_extension(content_type)
        if guessed_ext:
            filename += guessed_ext

    return filename


def _download_url_to_path(
    url: str,
    save_path: Optional[str] = None,
    output_dir: str = DEFAULT_ARTIFACT_DIR,
    timeout: int = 30,
) -> Dict[str, Any]:
    """Download a URL to a local path and return structured metadata."""
    try:
        import requests
    except ImportError as e:
        raise RuntimeError("requests not installed") from e

    url = _sanitize_url_input(url)
    normalized_url = _normalize_download_url(url)
    session = requests.Session()
    response = session.get(
        normalized_url,
        headers=_http_headers(),
        stream=True,
        timeout=timeout,
        allow_redirects=True,
    )
    response.raise_for_status()

    if not save_path:
        filename = _guess_filename_from_response(normalized_url, response)
        save_path = os.path.join(output_dir, filename)
    elif not os.path.isabs(save_path):
        save_path = os.path.join(output_dir, save_path)

    os.makedirs(os.path.dirname(save_path), exist_ok=True)

    with open(save_path, "wb") as f:
        for chunk in response.iter_content(chunk_size=8192):
            if chunk:
                f.write(chunk)

    return {
        "requested_url": url,
        "normalized_url": normalized_url,
        "final_url": response.url or normalized_url,
        "save_path": save_path,
        "size": os.path.getsize(save_path),
        "content_type": response.headers.get("Content-Type", "").split(";", 1)[0].strip().lower(),
    }


def _normalize_language_hint(language: Optional[str]) -> Optional[str]:
    """Normalize language hints to a short primary subtag such as 'en' or 'zh'."""
    if language is None:
        return None
    normalized = str(language).strip().replace("_", "-").lower()
    if not normalized:
        return None
    if "-" in normalized:
        normalized = normalized.split("-", 1)[0].strip()
    return normalized or None


def _stringify_cache_part(value: Any) -> str:
    """Convert cache-key parts into stable strings."""
    if value is None:
        return ""
    if isinstance(value, (dict, list)):
        try:
            return json.dumps(value, ensure_ascii=False, sort_keys=True)
        except Exception:
            return str(value)
    return str(value)


def _build_file_task_cache_key(file_path: str, *parts: Any) -> Tuple[str, int, int, Tuple[str, ...]]:
    """Build a per-process cache key tied to file identity plus task-specific parts."""
    abs_path = os.path.abspath(file_path)
    try:
        stat = os.stat(abs_path)
        size = int(stat.st_size)
        mtime_ns = int(stat.st_mtime_ns)
    except OSError:
        size = -1
        mtime_ns = -1
    normalized_parts = tuple(_stringify_cache_part(part).strip() for part in parts)
    return (abs_path, size, mtime_ns, normalized_parts)


def canonicalize_image_recognition_prompt_family(task: str, custom_prompt: str = "") -> str:
    """Group semantically similar vision prompts into stable task families for caching/loop control."""
    task_norm = str(task or "describe").strip().lower() or "describe"
    prompt_norm = re.sub(r"\s+", " ", str(custom_prompt or "").strip().lower())
    prompt_norm = re.sub(r"[\"'`]+", "", prompt_norm)
    if not prompt_norm:
        return f"{task_norm}:default"

    def has_any(*terms: str) -> bool:
        return any(term in prompt_norm for term in terms)

    def has_all(*terms: str) -> bool:
        return all(term in prompt_norm for term in terms)

    if task_norm in {"extract_numbers", "extract_by_color"}:
        colors = {color for color in ("red", "green", "purple") if color in prompt_norm}
        if colors == {"red", "green"}:
            return f"{task_norm}:color_group:red_green"
        if colors == {"red"}:
            return f"{task_norm}:color_group:red_only"
        if colors == {"green"}:
            return f"{task_norm}:color_group:green_only"
        if colors == {"purple"} and has_any("polygon", "side", "label", "vertices", "shape"):
            return f"{task_norm}:polygon_purple_labels"
        if has_all("numbers", "colors") or has_all("numbers", "their colors"):
            return f"{task_norm}:numbers_with_colors"

    if task_norm == "chess":
        if has_any("fen", "forsyth-edwards"):
            return "chess:fen"
        if has_any("next move", "best move", "guarantees a win", "algebraic notation"):
            return "chess:best_move"
        if has_any("board layout", "list all pieces", "piece positions", "coordinates", "side to move", "whose turn"):
            return "chess:board_state"
        return "chess:generic"

    if task_norm == "music_sheet":
        if has_any("bass clef", "staff", "line", "space", "note", "letter name"):
            return "music_sheet:notes_layout"
        return "music_sheet:generic"

    if task_norm == "geometry":
        if has_any("exact area", "calculate the exact area", "compute its area"):
            return "geometry:area_direct"
        if has_any("coordinate", "vertex", "vertices"):
            return "geometry:vertices"
        if has_any("side length", "side lengths", "labels", "clockwise", "purple number", "purple numbers"):
            return "geometry:side_measurements"
        if has_any("shape", "layout", "quadrilateral", "describe"):
            return "geometry:describe_shape"
        return "geometry:generic"

    if task_norm == "fractions":
        if has_any("classify", "quiz", "mixed number", "improper fraction"):
            return "fractions:classify"
        if has_any("comma-separated", "comma separated", "list all fractions", "sample problem", "answer to each"):
            return "fractions:list"
        return "fractions:generic"

    if task_norm == "extract_text":
        if has_any("python script", "indentation", "line breaks", "exactly as they appear", "exactly as it appears", "all characters"):
            return "extract_text:code_ocr"
        if "butterfat" in prompt_norm:
            return "extract_text:nutrition_butterfat"
        if has_all("saturated", "fat"):
            return "extract_text:nutrition_saturated_fat"
        if has_all("brand", "harness"):
            return "extract_text:brand_harness"
        return f"{task_norm}:raw:{prompt_norm}"

    return f"{task_norm}:raw:{prompt_norm}"


def _store_bounded_cache(cache: Dict[Any, Any], key: Any, value: Any, max_entries: int = 128) -> None:
    """Store into a small FIFO-ish cache without growing unbounded."""
    cache[key] = value
    if len(cache) > max_entries:
        oldest_key = next(iter(cache))
        if oldest_key != key:
            cache.pop(oldest_key, None)
        else:
            keys = list(cache.keys())
            if len(keys) > 1:
                cache.pop(keys[1], None)


def _normalize_lookup_name(value: str) -> str:
    """Normalize a filename-like string for fuzzy attachment matching."""
    return re.sub(r"[^a-z0-9]+", "", str(value or "").lower())


def _iter_attachment_files(attachments_dir: str) -> List[str]:
    """List attachment files under a directory, capped to avoid expensive scans."""
    if not attachments_dir or not os.path.isdir(attachments_dir):
        return []

    files: List[str] = []
    for root, _, names in os.walk(attachments_dir):
        for name in names:
            full_path = os.path.join(root, name)
            if os.path.isfile(full_path):
                files.append(full_path)
                if len(files) >= _MAX_ATTACHMENT_CANDIDATES:
                    return files
    return files


def _resolve_file_argument(file_path: Any, attachments_dir: str = "") -> Tuple[Any, Optional[str]]:
    """
    Resolve model-supplied file arguments against the working tree and attachments.

    Models often invent partial GAIA paths, omit the sample directory, or slightly
    mistype sample ids. This resolver keeps the public tool API unchanged while
    recovering the intended local attachment when there is a strong match.
    """
    if not isinstance(file_path, str) or not file_path.strip():
        return file_path, None

    original = file_path
    expanded = os.path.expanduser(file_path.strip().strip('"').strip("'"))
    candidates: List[str] = []

    def add_candidate(path_value: str) -> None:
        if path_value and path_value not in candidates:
            candidates.append(path_value)

    add_candidate(expanded)
    add_candidate(os.path.abspath(expanded))
    if attachments_dir:
        add_candidate(os.path.join(attachments_dir, expanded))
        add_candidate(os.path.join(attachments_dir, os.path.basename(expanded)))

    for candidate in candidates:
        if os.path.exists(candidate):
            resolved = os.path.abspath(candidate)
            note = None if resolved == os.path.abspath(expanded) else f"Resolved file_path: {original} -> {resolved}"
            return resolved, note

    attachment_files = _iter_attachment_files(attachments_dir)
    if not attachment_files:
        return original, None

    requested_name = os.path.basename(expanded)
    requested_stem = Path(requested_name).stem
    requested_suffix = Path(requested_name).suffix.lower()
    requested_norm = _normalize_lookup_name(requested_name)
    requested_stem_norm = _normalize_lookup_name(requested_stem)

    scored: List[Tuple[float, str]] = []
    for candidate in attachment_files:
        cand_name = os.path.basename(candidate)
        cand_stem = Path(cand_name).stem
        cand_suffix = Path(cand_name).suffix.lower()
        cand_norm = _normalize_lookup_name(cand_name)
        cand_stem_norm = _normalize_lookup_name(cand_stem)

        score = 0.0
        if requested_name and cand_name.lower() == requested_name.lower():
            score = 1.0
        elif requested_stem and cand_stem.lower() == requested_stem.lower():
            score = 0.96
        elif requested_norm and cand_norm == requested_norm:
            score = 0.94
        elif requested_stem_norm and cand_stem_norm == requested_stem_norm:
            score = 0.92
        elif (
            requested_stem_norm
            and cand_stem_norm
            and len(requested_stem_norm) >= 8
            and requested_stem_norm[:8] == cand_stem_norm[:8]
        ):
            score = 0.90
        else:
            ratio = difflib.SequenceMatcher(None, requested_norm, cand_norm).ratio() if requested_norm and cand_norm else 0.0
            stem_ratio = (
                difflib.SequenceMatcher(None, requested_stem_norm, cand_stem_norm).ratio()
                if requested_stem_norm and cand_stem_norm else 0.0
            )
            score = max(ratio, stem_ratio)
            if requested_stem_norm and len(requested_stem_norm) >= 8:
                if requested_stem_norm in cand_stem_norm or cand_stem_norm in requested_stem_norm:
                    score = max(score, 0.88)

        if requested_suffix and cand_suffix and requested_suffix != cand_suffix:
            score -= 0.08
        scored.append((score, candidate))

    scored.sort(reverse=True, key=lambda item: item[0])
    if scored and scored[0][0] >= 0.84:
        resolved = os.path.abspath(scored[0][1])
        return resolved, f"Resolved file_path: {original} -> {resolved}"

    return original, None


def _format_input_files_for_python(attachments_dir: str) -> Tuple[str, List[str]]:
    """Return attachment context for python_executor."""
    files = [os.path.abspath(p) for p in _iter_attachment_files(attachments_dir)]
    return (os.path.abspath(attachments_dir) if attachments_dir else "", files)


def _sanitize_url_input(url: Any) -> str:
    """Clean common malformed URL strings produced by models."""
    value = str(url or "").strip().strip('"').strip("'")
    value = re.sub(r"\s+", "", value)
    value = value.replace("hxxps://", "https://").replace("hxxp://", "http://")
    value = value.replace("ht./tps./:.///", "https://")
    value = value.replace("ht./tp./:.///", "http://")
    value = re.sub(r"^https?:/+(?=[A-Za-z0-9.-]+\.)", lambda m: m.group(0).split(":", 1)[0] + "://", value)
    value = value.replace("https:///", "https://").replace("http:///", "http://")
    return value


def _is_search_or_captcha_url(url: str) -> bool:
    """Identify URLs that should not be browsed directly as content pages."""
    parsed = urlparse(url)
    host = parsed.netloc.lower()
    path = parsed.path.lower()
    if "google." in host and (path.startswith("/search") or path.startswith("/sorry")):
        return True
    if "bing.com" in host and path.startswith("/search"):
        return True
    if "search.yahoo.com" in host:
        return True
    return False


# =============================================================================
# Active Tool 1: web_search
# =============================================================================

def _fetch_page_text_light(url: str, max_length: int = 2000) -> str:
    """Lightweight page fetch using requests only (for search result expansion)."""
    url = _sanitize_url_input(url)
    if _looks_like_pdf_url(url) or _is_search_or_captcha_url(url):
        return ""
    cache = getattr(_fetch_page_text_light, "_cache", {})
    cache_key = (url, max_length)
    cached = cache.get(cache_key)
    if isinstance(cached, str):
        return cached
    try:
        import requests
        from bs4 import BeautifulSoup
        resp = requests.get(_normalize_download_url(url), headers=_http_headers(), timeout=8)
        resp.raise_for_status()
        if "application/pdf" in resp.headers.get("Content-Type", "").lower():
            return ""
        soup = BeautifulSoup(resp.content, 'html.parser')
        for tag in soup(["script", "style", "noscript", "nav", "footer"]):
            tag.extract()
        text = soup.get_text()
        lines = (line.strip() for line in text.splitlines())
        text = '\n'.join(line for line in lines if line)
        if len(text) > max_length:
            text = text[:max_length] + "..."
        _store_bounded_cache(cache, cache_key, text, max_entries=256)
        _fetch_page_text_light._cache = cache
        return text
    except Exception:
        return ""


def web_search(query: str, engine: str = "duckduckgo", max_results: int = 5) -> str:
    """
    Search the web using DuckDuckGo.
    Returns: Title + Snippet for each result.
    Auto-expands the top 2 results by fetching their page content.
    """
    try:
        from ddgs import DDGS
    except ImportError:
        return "[ERROR] ddgs not installed. Run: pip install ddgs"

    query = str(query or "").strip()
    try:
        max_results = max(1, min(int(max_results), 8))
    except Exception:
        max_results = 5
    cache = getattr(web_search, "_cache", {})
    cache_key = (query.lower(), engine, max_results)
    cached = cache.get(cache_key)
    if isinstance(cached, str) and cached:
        return f"[CACHE HIT]\n{cached}"

    try:
        last_error = None
        results = []
        for backend in ("bing", "duckduckgo"):
            try:
                with DDGS() as ddgs:
                    results = list(ddgs.text(query, max_results=max_results, backend=backend))
                if results:
                    break
            except Exception as e:
                last_error = e

        if not results:
            if last_error is not None:
                return f"[ERROR] Web search failed: {str(last_error)}"
            return f"[INFO] No results found for query: {query}"

        output_lines = [f"Search results for: '{query}'\n"]
        expand_top = int(os.environ.get("WEB_SEARCH_EXPAND_TOP", "2") or 2)
        for i, r in enumerate(results, 1):
            title = r.get("title", "No title")
            snippet = r.get("body", "No snippet")
            url = _sanitize_url_input(r.get("href", ""))
            output_lines.append(f"[{i}] {title}")
            output_lines.append(f"    {snippet}")
            if url:
                output_lines.append(f"    URL: {url}")

            # Expand only the top few results. Over-expansion triggers rate limits.
            if i <= expand_top and url:
                page_text = _fetch_page_text_light(url, max_length=2000)
                if page_text:
                    output_lines.append(f"    --- Expanded Content ---")
                    output_lines.append(f"    {page_text[:2000]}")
                    output_lines.append(f"    --- End Expanded ---")

            output_lines.append("")

        output = "\n".join(output_lines)
        _store_bounded_cache(cache, cache_key, output, max_entries=128)
        web_search._cache = cache
        return output

    except Exception as e:
        return f"[ERROR] Web search failed: {str(e)}"


# =============================================================================
# Active Tool 2: excel_reader
# =============================================================================

def _normalize_column_label(value: Any) -> str:
    """Normalize table column labels for fuzzy lookup."""
    return re.sub(r"[^a-z0-9]+", "_", str(value or "").strip().lower()).strip("_")


def _query_terms(query: Optional[str]) -> List[str]:
    """Extract lightweight search terms from a natural-language table query."""
    if not query:
        return []
    raw_terms = re.findall(r"[A-Za-z0-9_][A-Za-z0-9_ -]{1,40}", str(query))
    terms: List[str] = []
    stop = {"what", "which", "where", "when", "how", "the", "and", "for", "with", "from", "column", "columns", "row", "rows"}
    for term in raw_terms:
        cleaned = term.strip(" -_").lower()
        if cleaned and cleaned not in stop and len(cleaned) > 1:
            terms.append(cleaned)
    return terms[:20]


def _relevant_columns(df, query: Optional[str], limit: int = 12) -> List[Any]:
    """Pick columns likely relevant to the query while preserving order."""
    if not query:
        return list(df.columns[:limit])

    normalized_cols = {col: _normalize_column_label(col) for col in df.columns}
    terms = _query_terms(query)
    scored: List[Tuple[float, Any]] = []
    for col, norm_col in normalized_cols.items():
        score = 0.0
        col_text = str(col).lower()
        for term in terms:
            norm_term = _normalize_column_label(term)
            if not norm_term:
                continue
            if norm_term == norm_col:
                score = max(score, 1.0)
            elif norm_term in norm_col or norm_col in norm_term:
                score = max(score, 0.86)
            else:
                score = max(score, difflib.SequenceMatcher(None, norm_term, norm_col).ratio())
            if term in col_text:
                score = max(score, 0.9)
        if score >= 0.55:
            scored.append((score, col))

    scored.sort(reverse=True, key=lambda item: item[0])
    selected = [col for _, col in scored[:limit]]
    if not selected:
        selected = list(df.columns[:limit])
    return selected


def _format_dataframe_inspection(df, file_path: str, query: Optional[str] = None, sheet_name: Optional[str] = None) -> str:
    """Return a compact but useful table inspection report."""
    output = []
    label = os.path.basename(file_path)
    if sheet_name:
        label += f" / sheet={sheet_name}"
    output.append(f"File: {label}")
    output.append(f"Shape: {df.shape[0]} rows x {df.shape[1]} columns")

    normalized_map = {str(col): _normalize_column_label(col) for col in df.columns}
    output.append(f"\nColumns ({len(df.columns)}): {list(df.columns)}")
    output.append(f"Normalized column map: {normalized_map}")

    profile_rows = []
    for col in df.columns:
        series = df[col]
        non_null = int(series.notna().sum())
        sample_values = [str(v) for v in series.dropna().astype(str).head(4).tolist()]
        unique_count = int(series.nunique(dropna=True))
        profile_rows.append({
            "column": str(col),
            "dtype": str(series.dtype),
            "non_null": non_null,
            "unique": unique_count,
            "sample_values": "; ".join(sample_values[:4]),
        })

    try:
        import pandas as pd
        profile_df = pd.DataFrame(profile_rows)
        output.append("\nColumn profile:")
        output.append(profile_df.to_markdown(index=False))
    except Exception:
        output.append(f"\nColumn profile: {profile_rows}")

    selected_columns = _relevant_columns(df, query)
    if query:
        output.append(f"\nQuery hint: {query}")
        output.append(f"Relevant columns by fuzzy match: {selected_columns}")

    preview_columns = selected_columns or list(df.columns[:12])
    output.append("\nPreview rows:")
    try:
        output.append(df.loc[:, preview_columns].head(10).to_markdown(index=False))
    except Exception:
        output.append(df.head(10).to_markdown(index=False))

    return "\n".join(output)


def excel_reader(file_path: str, sheet: Optional[str] = None, query: Optional[str] = None) -> str:
    """
    Read Excel/CSV files using pandas.
    Returns: table shape, normalized column map, dtypes, samples, and preview rows.
    """
    try:
        import pandas as pd
    except ImportError:
        return "[ERROR] pandas not installed. Run: pip install pandas openpyxl"

    if not os.path.exists(file_path):
        return f"[ERROR] File not found: {file_path}"

    try:
        ext = Path(file_path).suffix.lower()

        if ext == ".csv":
            df = pd.read_csv(file_path)
        elif ext in [".xlsx", ".xls"]:
            if sheet:
                df = pd.read_excel(file_path, sheet_name=sheet)
            else:
                # Read all sheets if no specific sheet requested
                excel_file = pd.ExcelFile(file_path)
                sheet_names = excel_file.sheet_names

                if len(sheet_names) == 1:
                    df = pd.read_excel(file_path, sheet_name=sheet_names[0])
                else:
                    # Multiple sheets - return info about all sheets
                    output = [f"Excel file contains {len(sheet_names)} sheets: {sheet_names}\n"]
                    for sn in sheet_names:
                        sheet_df = pd.read_excel(file_path, sheet_name=sn)
                        output.append(f"\n=== Sheet: {sn} ===")
                        output.append(_format_dataframe_inspection(sheet_df, file_path, query=query, sheet_name=sn))
                    return "\n".join(output)
        else:
            return f"[ERROR] Unsupported file format: {ext}"

        return _format_dataframe_inspection(df, file_path, query=query, sheet_name=sheet)

    except Exception as e:
        return f"[ERROR] Failed to read Excel file: {str(e)}"


# =============================================================================
# Active Tool 3: audio_transcription
# =============================================================================

def _transcribe_via_api(file_path: str, language: Optional[str] = None) -> Optional[str]:
    """
    Transcribe audio using API-based Whisper (OpenAI-compatible endpoint).

    Environment Variables:
        - ASR_MODEL_NAME: Primary ASR model (default: whisper-Breeze-ASR-25)
        - ASR_FALLBACK_MODEL: Fallback ASR model (default: whisper-large-v2)
        - ASR_API_BASE: API endpoint (default: LLM_API_BASE)
        - ASR_API_KEY: API key (default: LLM_API_KEY)
        - ASR_TIMEOUT: API timeout in seconds (default: 120)
        - ASR_MAX_RETRIES: Max retry attempts (default: 3)

    Returns:
        Transcription text if successful, None if API unavailable/failed.
    """
    import time
    import random

    # API configuration (similar to vision model pattern)
    api_base = normalize_api_base(os.environ.get("ASR_API_BASE") or get_api_base())
    api_key = os.environ.get("ASR_API_KEY") or get_api_key()

    if not api_key or not api_base:
        print("[INFO] ASR API not configured, will use local fallback")
        return None

    # Model configuration
    asr_model = os.environ.get("ASR_MODEL_NAME", "whisper-Breeze-ASR-25")
    fallback_model = os.environ.get("ASR_FALLBACK_MODEL", "whisper-large-v2")
    timeout_seconds = float(os.environ.get("ASR_TIMEOUT", "120"))
    max_retries = int(os.environ.get("ASR_MAX_RETRIES", "3"))

    # Prepare the transcription endpoint URL
    transcription_url = f"{api_base}/audio/transcriptions"

    # Try primary model first, then fallback
    models_to_try = [asr_model]
    if fallback_model and fallback_model != asr_model:
        models_to_try.append(fallback_model)

    for model_name in models_to_try:
        print(f"[INFO] Attempting API transcription with {model_name}...")

        for attempt in range(max_retries):
            try:
                import requests

                # Prepare multipart form data
                with open(file_path, "rb") as audio_file:
                    files = {"file": (os.path.basename(file_path), audio_file)}
                    data = {"model": model_name, "response_format": "json"}

                    # Add language hint if provided
                    normalized_lang = _normalize_language_hint(language)
                    if normalized_lang:
                        data["language"] = normalized_lang

                    headers = {"Authorization": f"Bearer {api_key}"}

                    response = requests.post(
                        transcription_url,
                        headers=headers,
                        files=files,
                        data=data,
                        timeout=timeout_seconds
                    )

                if response.status_code == 200:
                    result = response.json()
                    transcription = result.get("text", "")
                    detected_lang = result.get("language", "unknown")
                    print(f"[INFO] API transcription successful with {model_name}")
                    return f"Audio file: {os.path.basename(file_path)}\nDetected language: {detected_lang}\n\nTranscription:\n{transcription}"

                elif response.status_code == 429:
                    # Rate limited - exponential backoff
                    wait_time = (2 ** attempt) + random.uniform(0, 1)
                    print(f"[WARN] Rate limited, waiting {wait_time:.1f}s before retry...")
                    time.sleep(wait_time)
                    continue

                elif response.status_code >= 500:
                    # Server error - retry with backoff. Log a short body snippet
                    # because some upstream gateways return generic HTML errors
                    # that would otherwise be invisible in experiment logs.
                    body_preview = (response.text or "")[:200].replace("\n", " ")
                    wait_time = (2 ** attempt) + random.uniform(0, 1)
                    print(
                        f"[WARN] Server error {response.status_code} from ASR API "
                        f"(body: {body_preview!r}), waiting {wait_time:.1f}s..."
                    )
                    time.sleep(wait_time)
                    continue

                else:
                    print(f"[WARN] API returned {response.status_code}: {response.text[:200]}")
                    break  # Don't retry on client errors

            except requests.exceptions.Timeout:
                print(f"[WARN] API timeout (attempt {attempt + 1}/{max_retries})")
                continue
            except requests.exceptions.RequestException as e:
                print(f"[WARN] API request failed: {str(e)}")
                break
            except Exception as e:
                print(f"[WARN] Unexpected error in API transcription: {str(e)}")
                break

        print(f"[WARN] Failed with {model_name}, trying next model...")

    print("[INFO] All API models failed, falling back to local Whisper")
    return None


def _transcribe_via_local(file_path: str, language: Optional[str] = None) -> str:
    """
    Transcribe audio using local OpenAI Whisper model (lazy loaded).

    This is the fallback when API is unavailable.
    """
    global _whisper_model

    try:
        import whisper
    except ImportError:
        return "[ERROR] openai-whisper not installed. Run: pip install openai-whisper"

    try:
        # Lazy load model (base model for speed/accuracy balance)
        if _whisper_model is None:
            import torch
            best_device = "cpu"
            if torch.cuda.is_available():
                max_free_mem = 0
                best_gpu_idx = 0
                for i in range(torch.cuda.device_count()):
                    try:
                        free_mem, _ = torch.cuda.mem_get_info(i)
                        if free_mem > max_free_mem:
                            max_free_mem = free_mem
                            best_gpu_idx = i
                    except Exception:
                        pass

                # Require at least ~2GB free for whisper base
                if max_free_mem > 2 * 1024 * 1024 * 1024:
                    best_device = f"cuda:{best_gpu_idx}"
                    print(f"[INFO] Selected {best_device} for Whisper (Free: {max_free_mem / 1024**3:.2f} GB)")
                else:
                    print(f"[WARN] Not enough VRAM on any GPU. Falling back to CPU for Whisper.")

            print(f"[INFO] Loading local Whisper model (base) on {best_device}...")
            _whisper_model = whisper.load_model("base", device=best_device)

        # Transcribe options
        options = {}
        normalized_language = _normalize_language_hint(language)
        if normalized_language:
            if normalized_language != str(language).strip():
                print(f"[INFO] Normalized language '{language}' -> '{normalized_language}'")
            options["language"] = normalized_language

        # FP16 inference for speed (if CUDA available)
        import torch
        if torch.cuda.is_available():
            options["fp16"] = True

        print(f"[INFO] Transcribing with local Whisper: {os.path.basename(file_path)}...")
        result = _whisper_model.transcribe(file_path, **options)
        print(f"[INFO] Local transcription completed")

        transcription = result.get("text", "")
        output = [
            f"Audio file: {os.path.basename(file_path)}",
            f"Detected language: {result.get('language', 'unknown')}",
            f"\nTranscription:",
            transcription
        ]
        return "\n".join(output)

    except Exception as e:
        return f"[ERROR] Local audio transcription failed: {str(e)}"


def audio_transcription(file_path: str, language: Optional[str] = None) -> str:
    """
    Transcribe audio using Whisper with API-first strategy and local fallback.

    This tool uses a two-tier architecture similar to image_recognition:
    1. Primary: API-based whisper-Breeze-ASR-25
    2. Fallback: Local whisper base model (faster, ~10% WER)

    Args:
        file_path: Path to the audio file (mp3, wav, m4a, etc.)
        language: Optional language hint (e.g., "en", "zh", "ja")

    Environment Variables:
        - ASR_MODEL_NAME: Primary ASR model (default: whisper-Breeze-ASR-25)
        - ASR_FALLBACK_MODEL: Fallback ASR model (default: whisper-large-v2)
        - ASR_API_BASE: API endpoint (default: LLM_API_BASE)
        - ASR_API_KEY: API key (default: LLM_API_KEY)
        - ASR_USE_LOCAL_ONLY: Set to "true" to skip API and use local only

    Returns:
        Transcription result with detected language and text.
    """
    if not os.path.exists(file_path):
        return f"[ERROR] Audio file not found: {file_path}"

    normalized_language = _normalize_language_hint(language)
    cache_key = _build_file_task_cache_key(file_path, normalized_language or "")
    result_cache = getattr(audio_transcription, "_result_cache", {})
    cached_result = result_cache.get(cache_key)
    if isinstance(cached_result, str) and cached_result:
        print(f"[INFO] Returning cached audio transcription for {os.path.basename(file_path)}")
        return cached_result

    # Check if local-only mode is enabled
    use_local_only = os.environ.get("ASR_USE_LOCAL_ONLY", "false").lower() == "true"

    if not use_local_only:
        # Try API-based transcription first
        api_result = _transcribe_via_api(file_path, normalized_language)
        if api_result is not None:
            _store_bounded_cache(result_cache, cache_key, api_result, max_entries=128)
            audio_transcription._result_cache = result_cache
            return api_result

    # Fallback to local Whisper
    local_result = _transcribe_via_local(file_path, normalized_language)
    if isinstance(local_result, str) and local_result and not local_result.startswith("[ERROR]"):
        _store_bounded_cache(result_cache, cache_key, local_result, max_entries=128)
        audio_transcription._result_cache = result_cache
    return local_result


# =============================================================================
# Active Tool 4: file_reader
# =============================================================================

def file_reader(file_path: str, **kwargs) -> str:
    """
    Read various file formats:
    - .docx: Use python-docx, iterate paragraphs AND tables
    - .json/.jsonld/.xml/.txt/.py/.pdb/.csv-like text: Read as text
    
    Args:
        file_path: Path to the file
        **kwargs: Extra arguments (e.g., 'page') are ignored
    """
    if not os.path.exists(file_path):
        return f"[ERROR] File not found: {file_path}"

    ext = Path(file_path).suffix.lower()

    try:
        # DOCX files
        if ext == ".docx":
            try:
                from docx import Document
            except ImportError:
                return "[ERROR] python-docx not installed. Run: pip install python-docx"

            doc = Document(file_path)
            output = [f"Document: {os.path.basename(file_path)}\n"]

            # Extract paragraphs
            output.append("=== Paragraphs ===")
            for i, para in enumerate(doc.paragraphs, 1):
                if para.text.strip():
                    output.append(f"{para.text}")

            # Extract tables
            if doc.tables:
                output.append(f"\n=== Tables ({len(doc.tables)}) ===")
                for t_idx, table in enumerate(doc.tables, 1):
                    output.append(f"\n--- Table {t_idx} ---")
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        output.append(row_text)

            return "\n".join(output)

        # JSON/JSONLD files
        elif ext in [".json", ".jsonld"]:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            try:
                data = json.loads(content)
                return f"JSON content ({os.path.basename(file_path)}):\n{json.dumps(data, indent=2, ensure_ascii=False)}"
            except json.JSONDecodeError:
                return f"Raw content ({os.path.basename(file_path)}):\n{content}"

        # XML files
        elif ext == ".xml":
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            return f"XML content ({os.path.basename(file_path)}):\n{content}"

        # Plain-text-like files
        elif ext in {".txt", ".py", ".md", ".csv", ".tsv", ".pdb", ".sdf", ".mol", ".cpp", ".c", ".h"}:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            return f"Text content ({os.path.basename(file_path)}):\n{content}"

        else:
            return f"[ERROR] Unsupported file format: {ext}. Supported: .docx, .json, .jsonld, .xml, .txt, .py, .pdb, and common text/code files"

    except Exception as e:
        return f"[ERROR] Failed to read file: {str(e)}"



# =============================================================================
# Active Tool 5: python_executor
# =============================================================================

def _normalize_python_code(code: Any) -> str:
    """Repair common escaped-code formatting without changing normal code."""
    text = str(code or "")
    stripped = text.strip()
    if (
        len(stripped) >= 2
        and stripped[0] == stripped[-1]
        and stripped[0] in {"'", '"'}
        and "\\n" in stripped
    ):
        try:
            parsed = json.loads(stripped) if stripped[0] == '"' else None
            if isinstance(parsed, str):
                return parsed
        except Exception:
            pass
        stripped = stripped[1:-1]

    if "\\n" in stripped and stripped.count("\n") <= 1:
        code_like = re.search(r"(import|from|print|for |if |def |class |with |pd\.|open\(|=)", stripped)
        if code_like:
            stripped = stripped.replace("\\n", "\n").replace("\\t", "\t")
    return stripped


def python_executor(
    code: str,
    output_dir: str = DEFAULT_ARTIFACT_DIR,
    attachments_dir: str = "",
    input_files: Optional[List[str]] = None,
    timeout_seconds: int = 30,
) -> str:
    """
    Execute Python code with captured stdout.
    Injects pandas (pd), numpy (np), and Biopython (Bio) into the execution namespace.
    """
    code = _normalize_python_code(code)
    lowered_code = code.lower()
    inspect_first_hint = (
        "[HINT] For files, tables, or spreadsheets, first run short inspection code "
        "(for example: print(df.columns), print(df.head()), print(df.shape)) before "
        "writing larger transformation scripts. Use ATTACHMENTS_DIR and INPUT_FILES "
        "instead of guessing GAIA attachment paths."
    )
    brute_force_markers = [
        "itertools.permutations(",
        "from itertools import permutations",
        " import permutations",
        " permutations(",
    ]
    if any(marker in lowered_code for marker in brute_force_markers):
        return (
            "[ERROR] Execution blocked: obvious brute-force permutation search is disabled "
            "in python_executor. Use a targeted algorithm, filtering, counting, or reasoning "
            "instead of exhaustive permutations."
        )

    exec_globals = {
        "__builtins__": __builtins__,
    }

    resolved_attachments_dir, resolved_input_files = _format_input_files_for_python(attachments_dir)
    if input_files:
        resolved_input_files = [os.path.abspath(str(p)) for p in input_files]
    exec_globals["ATTACHMENTS_DIR"] = resolved_attachments_dir
    exec_globals["INPUT_FILES"] = resolved_input_files
    exec_globals["INPUT_FILE_NAMES"] = [os.path.basename(p) for p in resolved_input_files]

    # Add math and statistics modules
    import math
    import statistics
    exec_globals["math"] = math
    exec_globals["statistics"] = statistics
    # Add commonly used math/stats functions
    exec_globals["sqrt"] = math.sqrt
    exec_globals["log"] = math.log
    exec_globals["sin"] = math.sin
    exec_globals["cos"] = math.cos
    exec_globals["pi"] = math.pi
    exec_globals["e"] = math.e
    exec_globals["stdev"] = statistics.stdev
    exec_globals["pstdev"] = statistics.pstdev
    exec_globals["mean"] = statistics.mean
    exec_globals["median"] = statistics.median

    # Check Pandas
    libraries_status = []

    # Check Pandas
    try:
        import pandas as pd
        exec_globals["pd"] = pd
        exec_globals["pandas"] = pd
    except ImportError:
        libraries_status.append("pandas (MISSING)")

    # Check Numpy
    try:
        import numpy as np
        exec_globals["np"] = np
        exec_globals["numpy"] = np
    except ImportError:
        libraries_status.append("numpy (MISSING)")

    # Check Biopython
    try:
        import Bio
        exec_globals["Bio"] = Bio
        
        from Bio.PDB import PDBParser
        exec_globals["PDBParser"] = PDBParser
    except ImportError:
        libraries_status.append("biopython (MISSING)")

    try:
        import pypdf
        exec_globals["pypdf"] = pypdf
        exec_globals["PdfReader"] = pypdf.PdfReader
        sys.modules.setdefault("PyPDF2", pypdf)
    except ImportError:
        pass

    # If core libraries are missing, return error
    if "pandas (MISSING)" in libraries_status or "numpy (MISSING)" in libraries_status:
        return f"[SYSTEM ERROR] Essential libraries missing. Status: {', '.join(libraries_status)}. Please ask admin to install them."

    # Check Matplotlib
    plt = None
    try:
        import matplotlib
        matplotlib.use('Agg')  # Set to non-interactive mode
        import matplotlib.pyplot as plt
        exec_globals["plt"] = plt
        exec_globals["matplotlib"] = matplotlib
    except ImportError:
        pass

    # Set output directory
    if output_dir:
        exec_globals["OUTPUT_DIR"] = output_dir
        os.makedirs(output_dir, exist_ok=True)

    # --- Execution Phase ---
    
    # Capture stdout
    old_stdout = sys.stdout
    sys.stdout = captured_output = io.StringIO()

    result_value = None
    error_msg = None
    saved_files = []

    alarm_enabled = False
    old_alarm_handler = None

    def _timeout_handler(signum, frame):
        raise TimeoutError(f"python_executor exceeded {timeout_seconds} seconds")

    try:
        if timeout_seconds and hasattr(signal, "SIGALRM"):
            try:
                old_alarm_handler = signal.signal(signal.SIGALRM, _timeout_handler)
                signal.alarm(int(timeout_seconds))
                alarm_enabled = True
            except Exception:
                alarm_enabled = False

        # Execute the code
        exec(code, exec_globals)

        # Check for common result variable names
        for var_name in ['result', 'answer', 'output', 'value']:
            if var_name in exec_globals:
                result_value = exec_globals[var_name]
                break

        # Auto-save matplotlib figures if any are open
        if plt is not None and output_dir:
            fig_nums = plt.get_fignums()
            for i, fig_num in enumerate(fig_nums):
                fig = plt.figure(fig_num)
                fig_path = os.path.join(output_dir, f"plot_{i+1}.png")
                fig.savefig(fig_path, dpi=150, bbox_inches='tight')
                saved_files.append(fig_path)
            plt.close('all')

    except Exception as e:
        # Capture specific error types, help model debug
        error_msg = f"[ERROR] Execution failed: {type(e).__name__}: {str(e)}"
    except SystemExit as e:
        error_msg = f"[ERROR] Execution failed: Script called exit() with code {e.code}. Do not use exit()."

    finally:
        if alarm_enabled:
            try:
                signal.alarm(0)
                if old_alarm_handler is not None:
                    signal.signal(signal.SIGALRM, old_alarm_handler)
            except Exception:
                pass
        sys.stdout = old_stdout

    # --- Build return message ---
    output_lines = []
    stdout_content = captured_output.getvalue()

    if stdout_content:
        output_lines.append("=== Output ===")
        output_lines.append(stdout_content.strip())

    if result_value is not None:
        output_lines.append(f"\n=== Result ===")
        output_lines.append(str(result_value))

    if saved_files:
        output_lines.append(f"\n=== Generated Files (Use these paths for next steps) ===")
        for fp in saved_files:
            output_lines.append(f"  {fp}")

    if error_msg:
        output_lines.append(error_msg)
        if any(name in error_msg for name in ("KeyError", "IndexError", "ValueError", "SyntaxError", "NameError", "AttributeError", "ParserError")):
            output_lines.append(inspect_first_hint)
        if "ModuleNotFoundError" in error_msg or "Missing optional dependency" in error_msg or "ImportError" in error_msg:
            output_lines.append(
                "[HINT] Prefer built-in injected libraries first: pd, np, pypdf/PdfReader, math, statistics, Bio/PDBParser. "
                "If a missing optional package is essential, install it in the llm_planning environment."
            )

    if not output_lines:
        output_lines.append("[INFO] Code executed successfully (no output)")

    if resolved_input_files:
        output_lines.append("\n=== Available Input Files ===")
        for path in resolved_input_files[:25]:
            output_lines.append(f"- {path}")
        if len(resolved_input_files) > 25:
            output_lines.append(f"... ({len(resolved_input_files) - 25} more)")

    return "\n".join(output_lines)


# =============================================================================
# Active Tool 6: zip_extractor 
# =============================================================================

def zip_extractor(file_path: str, extract_to: Optional[str] = None, output_dir: str = DEFAULT_ARTIFACT_DIR) -> str:
    """
    Extract ZIP archive contents.
    Returns list of extracted file paths.
    """
    if not os.path.exists(file_path):
        return f"[ERROR] File not found: {file_path}"

    try:
        zip_name = Path(file_path).stem

        # Determine extraction directory
        if extract_to:
            extract_dir = extract_to
        elif output_dir:
            extract_dir = os.path.join(output_dir, f"{zip_name}_extracted")
        else:
            extract_dir = f"{zip_name}_extracted"

        os.makedirs(extract_dir, exist_ok=True)

        extracted_paths = []
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            
            zip_ref.extractall(extract_dir)
            
            
            for f in zip_ref.namelist():
                full_path = os.path.join(extract_dir, f)
                
                if os.path.isfile(full_path):
                    extracted_paths.append(full_path)

        # Build output
        output = [f"[SUCCESS] Extracted '{os.path.basename(file_path)}'"]
        output.append(f"Directory: {extract_dir}")
        
        
        output.append(f"\n=== Extracted File Paths (Use these for file_reader) ===")
        if not extracted_paths:
             output.append("(No files found in zip)")
        
        for p in extracted_paths:
            size = os.path.getsize(p)
            output.append(f"- {p} ({size} bytes)")

        return "\n".join(output)

    except zipfile.BadZipFile:
        return f"[ERROR] Invalid or corrupted ZIP file: {file_path}"
    except Exception as e:
        return f"[ERROR] Failed to extract ZIP: {str(e)}"


# =============================================================================
# Active Tool 7: pptx_reader
# =============================================================================

def pptx_reader(file_path: str, slide_number: Optional[int] = None) -> str:
    """
    Read PowerPoint presentations using python-pptx.
    Inserts [Slide N] markers for each slide.
    """
    if not os.path.exists(file_path):
        return f"[ERROR] File not found: {file_path}"

    try:
        from pptx import Presentation
    except ImportError:
        return "[ERROR] python-pptx not installed. Run: pip install python-pptx"

    try:
        prs = Presentation(file_path)
        total_slides = len(prs.slides)

        output = [f"PowerPoint: {os.path.basename(file_path)}"]
        output.append(f"Total slides: {total_slides}\n")

        for slide_idx, slide in enumerate(prs.slides, 1):
            # Skip if specific slide requested and this isn't it
            if slide_number is not None and slide_idx != slide_number:
                continue

            output.append(f"[Slide {slide_idx}]")

            # Extract text from all shapes
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            if slide_text:
                output.append("\n".join(slide_text))
            else:
                output.append("(No text content)")

            output.append("")  # Blank line between slides

        return "\n".join(output)

    except Exception as e:
        return f"[ERROR] Failed to read PowerPoint: {str(e)}"


# =============================================================================
# Active Tool 8: pdf_reader
# =============================================================================

def pdf_reader(file_path: str, page: Optional[int] = None) -> str:
    """
    Read PDF files using pypdf.
    Extracts text page by page.
    """
    source_label = file_path
    local_path = file_path

    if _is_http_url(file_path):
        try:
            downloaded = _download_url_to_path(file_path, output_dir=DEFAULT_ARTIFACT_DIR)
            local_path = downloaded["save_path"]
            source_label = f"{file_path}\nDownloaded to: {local_path}"
        except Exception as e:
            return f"[ERROR] Failed to fetch remote PDF: {str(e)}"

    if not os.path.exists(local_path):
        return f"[ERROR] File not found: {local_path}"

    try:
        from pypdf import PdfReader
    except ImportError:
        return "[ERROR] pypdf not installed. Run: pip install pypdf"

    try:
        reader = PdfReader(local_path)
        total_pages = len(reader.pages)

        output = [f"PDF: {os.path.basename(local_path)}"]
        if source_label != local_path:
            output.append(f"Source: {source_label}")
        output.append(f"Total pages: {total_pages}\n")

        for page_idx, pdf_page in enumerate(reader.pages, 1):
            # Skip if specific page requested and this isn't it
            if page is not None and page_idx != page:
                continue

            output.append(f"[Page {page_idx}]")
            text = pdf_page.extract_text() or "(No text content)"
            output.append(text.strip())
            output.append("")  # Blank line between pages

        return "\n".join(output)

    except Exception as e:
        return f"[ERROR] Failed to read PDF: {str(e)}"





# =============================================================================
# Optional Tools (Implemented, may require extra deps)
# =============================================================================

# =============================================================================
# Active Tool 10: web_browser
# =============================================================================

def _extract_text_from_html(html: str, max_length: int = 8000) -> str:
    """Extract clean text from HTML content."""
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(html, 'html.parser')
    for tag in soup(["script", "style", "noscript", "nav", "footer", "header"]):
        tag.extract()
    text = soup.get_text()
    lines = (line.strip() for line in text.splitlines())
    chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
    text = '\n'.join(chunk for chunk in chunks if chunk)
    if len(text) > max_length:
        text = text[:max_length] + f"\n\n... (truncated, total {len(text)} chars)"
    return text


def _web_browser_playwright(url: str, timeout_ms: int = 15000) -> str:
    """Fetch page using Playwright headless browser (handles JS-rendered pages)."""
    from playwright.sync_api import sync_playwright
    with sync_playwright() as p:
        browser = p.chromium.launch(
            headless=True,
            args=[
                '--disable-blink-features=AutomationControlled',  # Hide automation
                '--disable-dev-shm-usage',
                '--no-sandbox',
            ]
        )
        context = browser.new_context(
            user_agent="Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
            viewport={"width": 1920, "height": 1080},
            locale="en-US",
            timezone_id="America/New_York",
        )

        # Set extra HTTP headers to appear more like a real browser
        context.set_extra_http_headers({
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,*/*;q=0.8",
            "Accept-Language": "en-US,en;q=0.9",
            "Accept-Encoding": "gzip, deflate, br",
            "DNT": "1",
            "Connection": "keep-alive",
            "Upgrade-Insecure-Requests": "1",
        })

        page = context.new_page()

        # Hide webdriver property
        page.add_init_script("""
            Object.defineProperty(navigator, 'webdriver', {
                get: () => undefined
            });
        """)

        page.goto(url, wait_until="domcontentloaded", timeout=timeout_ms)
        try:
            page.wait_for_load_state("networkidle", timeout=min(timeout_ms, 5000))
        except Exception:
            # Some pages never reach networkidle; DOM content is often enough.
            pass

        # Auto-click common consent/accept buttons (for sites like USGS, cookie banners)
        consent_selectors = [
            "button:has-text('Accept')", "button:has-text('Agree')", "button:has-text('OK')",
            "button:has-text('I Agree')", "button:has-text('Got it')", "button:has-text('Continue')",
            "a:has-text('Accept')", "a:has-text('Agree')", "a:has-text('OK')",
            "[class*='accept']", "[class*='consent']", "[id*='accept']", "[id*='consent']",
        ]
        for selector in consent_selectors:
            try:
                btn = page.locator(selector).first
                if btn.is_visible(timeout=500):
                    btn.click(timeout=1000)
                    try:
                        page.wait_for_load_state("networkidle", timeout=3000)
                    except Exception:
                        pass
                    break
            except:
                continue

        html = page.content()
        browser.close()
    return html


def _web_browser_requests(url: str) -> tuple[str, str, str]:
    """Fetch page using requests (fallback for static pages)."""
    import requests
    normalized_url = _normalize_download_url(url)
    response = requests.get(
        normalized_url,
        headers=_http_headers(accept="text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8"),
        timeout=12,
        allow_redirects=True,
    )
    response.raise_for_status()
    content_type = response.headers.get("Content-Type", "").split(";", 1)[0].strip().lower()
    return response.text, content_type, response.url or normalized_url


def _browse_pdf_url(url: str) -> str:
    """Handle PDF-like URLs through the existing download and PDF reader stack."""
    try:
        downloaded = _download_url_to_path(url, output_dir=DEFAULT_ARTIFACT_DIR)
    except Exception as e:
        return f"[ERROR] Failed to fetch PDF-like URL: {str(e)}"

    pdf_text = pdf_reader(downloaded["save_path"])
    return "\n".join([
        f"URL: {url}",
        f"Resolved URL: {downloaded['final_url']}",
        f"Downloaded to: {downloaded['save_path']}",
        "",
        pdf_text,
    ])


def web_browser(url: str, action: str = "navigate") -> str:
    """
    Web browser with Playwright headless (JS-rendering) + requests fallback.
    Fetches URL and extracts text content.

    Status: IMPLEMENTED
    Dependencies: playwright (preferred), requests + beautifulsoup4 (fallback)
    Install: pip install playwright && playwright install chromium
             pip install requests beautifulsoup4
    """
    try:
        import bs4  # noqa: F401 - needed by _extract_text_from_html
    except ImportError:
        return "[ERROR] beautifulsoup4 not installed. Run: pip install beautifulsoup4"

    original_url = str(url or "")
    sanitized_url = _sanitize_url_input(original_url)
    if _is_search_or_captcha_url(sanitized_url):
        return (
            "[ERROR] Web browser was asked to open a search/captcha page rather than a content page. "
            "Use web_search for search queries, then open a concrete result URL."
        )

    normalized_url = _normalize_download_url(sanitized_url)
    if _looks_like_pdf_url(normalized_url):
        return _browse_pdf_url(normalized_url)

    cache = getattr(web_browser, "_cache", {})
    cache_key = (normalized_url, action)
    cached = cache.get(cache_key)
    if isinstance(cached, str) and cached:
        return f"[CACHE HIT]\n{cached}"

    html = None
    method = "unknown"
    final_url = normalized_url

    # Try Playwright first (handles JS-rendered pages)
    try:
        html = _web_browser_playwright(normalized_url)
        method = "playwright"
    except ImportError:
        pass  # Playwright not installed, fall through to requests
    except Exception as e:
        print(f"  [WARN] Playwright failed for {normalized_url}: {e}, falling back to requests")

    # Fallback to requests
    if html is None:
        try:
            html, content_type, final_url = _web_browser_requests(normalized_url)
            if "application/pdf" in content_type or _looks_like_pdf_url(final_url):
                return _browse_pdf_url(final_url)
            method = "requests"
        except Exception as e:
            return f"[ERROR] Web browser failed: {str(e)}"

    try:
        text = _extract_text_from_html(html, max_length=8000)
        output = [
            f"URL: {original_url}",
            f"Sanitized URL: {sanitized_url}",
            f"Resolved URL: {final_url}",
            f"Action: {action}",
            f"Method: {method}",
            f"\nExtracted Text:\n",
            text
        ]
        result = "\n".join(output)
        _store_bounded_cache(cache, cache_key, result, max_entries=128)
        web_browser._cache = cache
        return result
    except Exception as e:
        return f"[ERROR] Web browser text extraction failed: {str(e)}"

# =============================================================================
# Active Tool 11: calculator
# =============================================================================

def calculator(expression: str) -> str:
    """
    Calculator that delegates to python_executor.
    Wraps expression in print() statement for evaluation.

    Supports: Basic math (+, -, *, /, **, %), math functions (sqrt, log, sin, cos, pi, e),
    and statistics functions (stdev, pstdev, mean, median).

    Examples:
        - "2 + 3 * 4"
        - "sqrt(16) + log(10)"
        - "stdev([1, 2, 3, 4, 5])"
        - "(stdev([1,2,3]) + pstdev([4,5,6])) / 2"

    Status: IMPLEMENTED (via python_executor)
    Dependencies: None (uses python_executor internally)
    """
    try:
        expression = expression.strip()
        code = f"result = {expression}\nprint(f'Result: {expression} = {{result}}')"
        result = python_executor(code, output_dir="")
        return result
    except Exception as e:
        return f"[ERROR] Calculator error: {str(e)}"




# =============================================================================
# Optional Tools / Mock Tools
# =============================================================================

def _save_vision_debug_record(
    output_dir: str,
    file_path: str,
    task: str,
    payload: Dict[str, Any],
) -> Optional[str]:
    """Persist one vision-debug payload under the run output directory."""
    if not output_dir:
        return None
    try:
        import time

        debug_dir = os.path.join(output_dir, "vision_debug")
        os.makedirs(debug_dir, exist_ok=True)
        stem = Path(file_path).stem or "image"
        safe_task = "".join(ch if ch.isalnum() or ch in {"-", "_"} else "_" for ch in task) or "task"
        save_path = os.path.join(debug_dir, f"vision_debug_{stem}_{safe_task}_{int(time.time() * 1000)}.json")
        with open(save_path, "w", encoding="utf-8") as f:
            json.dump(payload, f, indent=2, ensure_ascii=False)
        return save_path
    except Exception as e:
        print(f"  [WARN] Failed to save vision debug record: {e}")
        return None


def _classify_color_name(rgb: tuple[float, float, float]) -> str:
    """Map an RGB triple to a coarse color name."""
    r, g, b = rgb
    if r > 185 and g < 140 and b < 140:
        return "red"
    if g > 150 and r < 170 and b < 150:
        return "green"
    if b > 160 and r < 150 and g < 170:
        return "blue"
    if r > 190 and g > 170 and b < 130:
        return "yellow"
    if r > 200 and 110 <= g <= 190 and b < 130:
        return "orange"
    if r > 160 and b > 160 and g < 150:
        return "purple"
    if r < 80 and g < 80 and b < 80:
        return "black"
    if r > 215 and g > 215 and b > 215:
        return "white"
    return f"rgb({int(r)},{int(g)},{int(b)})"


def _position_name(cx: float, cy: float, width: int, height: int) -> str:
    """Return a coarse position label."""
    if width <= 0 or height <= 0:
        return "unknown"
    horiz = "left" if cx < width / 3 else "right" if cx > 2 * width / 3 else "center"
    vert = "top" if cy < height / 3 else "bottom" if cy > 2 * height / 3 else "middle"
    if horiz == "center" and vert == "middle":
        return "center"
    if horiz == "center":
        return vert
    if vert == "middle":
        return horiz
    return f"{vert}-{horiz}"


def _run_tesseract_makebox(image_path: str, psm: int, image_height: int) -> List[Dict[str, Any]]:
    """Run tesseract in makebox mode and return digit-level boxes."""
    cmd = [
        _find_tesseract_binary(),
        image_path,
        "stdout",
        "--psm",
        str(psm),
        "-c",
        "tessedit_char_whitelist=0123456789.,-",
        "makebox",
    ]
    proc = subprocess.run(cmd, check=False, capture_output=True, text=True)
    if proc.returncode != 0:
        raise RuntimeError(proc.stderr.strip() or proc.stdout.strip() or "tesseract failed")
    boxes: List[Dict[str, Any]] = []
    for line in proc.stdout.splitlines():
        parts = line.strip().split()
        if len(parts) < 5:
            continue
        ch = parts[0]
        if ch not in "0123456789":
            continue
        left, bottom, right, top = map(int, parts[1:5])
        boxes.append({
            "char": ch,
            "left": left,
            "right": right,
            "top": image_height - top,
            "bottom": image_height - bottom,
            "width": max(1, right - left),
            "height": max(1, top - bottom),
        })
    return boxes


def _extract_colored_number_candidates(file_path: str) -> List[Dict[str, Any]]:
    """
    Use local OCR + color sampling for number/color-heavy images.

    If it cannot recover enough evidence, callers should fall back to a vision
    LLM rather than trusting a partial local parse.
    """
    import numpy as np
    import statistics
    from PIL import Image, ImageOps, ImageEnhance

    with Image.open(file_path) as raw_img:
        rgb_img = raw_img.convert("RGB")
        scale = 2
        rgb_large = rgb_img.resize(
            (rgb_img.width * scale, rgb_img.height * scale),
            Image.Resampling.LANCZOS,
        )

    gray = ImageOps.grayscale(rgb_large)
    boosted = ImageEnhance.Contrast(ImageOps.autocontrast(gray)).enhance(2.0)
    binary = boosted.point(lambda p: 255 if p > 160 else 0)
    variants = [
        ("gray_psm11", boosted, 11),
        ("binary_psm11", binary, 11),
        ("gray_psm6", boosted, 6),
    ]

    candidates: Dict[tuple[str, int, int], Dict[str, Any]] = {}
    for variant_name, variant_img, psm in variants:
        tmp_path = None
        try:
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                tmp_path = tmp.name
            variant_img.save(tmp_path)
            boxes = _run_tesseract_makebox(tmp_path, psm, variant_img.height)
            if not boxes:
                continue

            boxes.sort(key=lambda box: (box["top"], box["left"]))
            median_height = statistics.median(box["height"] for box in boxes)
            row_tol = max(18.0, median_height * 0.8)
            rows: List[List[Dict[str, Any]]] = []
            for box in boxes:
                center_y = (box["top"] + box["bottom"]) / 2
                placed = False
                for row_boxes in rows:
                    row_center = statistics.mean((b["top"] + b["bottom"]) / 2 for b in row_boxes)
                    if abs(center_y - row_center) <= row_tol:
                        row_boxes.append(box)
                        placed = True
                        break
                if not placed:
                    rows.append([box])

            for row_boxes in rows:
                row_boxes.sort(key=lambda box: box["left"])
                median_width = statistics.median(box["width"] for box in row_boxes)
                gap_threshold = max(18.0, median_width * 0.75)
                token_boxes: List[List[Dict[str, Any]]] = []
                current: List[Dict[str, Any]] = []
                prev_right: Optional[int] = None
                for box in row_boxes:
                    if not current:
                        current = [box]
                        prev_right = box["right"]
                        continue
                    gap = box["left"] - (prev_right or box["left"])
                    if gap > gap_threshold:
                        token_boxes.append(current)
                        current = [box]
                    else:
                        current.append(box)
                    prev_right = box["right"]
                if current:
                    token_boxes.append(current)

                for token in token_boxes:
                    text = "".join(box["char"] for box in token).strip()
                    if not text or not any(ch.isdigit() for ch in text):
                        continue
                    left = min(box["left"] for box in token)
                    top = min(box["top"] for box in token)
                    right = max(box["right"] for box in token)
                    bottom = max(box["bottom"] for box in token)
                    width = max(1, right - left)
                    height = max(1, bottom - top)

                    if len(text) > 3:
                        # Over-merged OCR chunks are not trustworthy enough for
                        # structured extraction; let the VLM fallback handle them.
                        continue

                    conf = 55.0 - max(0, len(text) - 2) * 10.0
                    if len(text) == 1:
                        conf -= 10.0

                    right = min(rgb_large.width, right)
                    bottom = min(rgb_large.height, bottom)
                crop = np.asarray(rgb_large.crop((left, top, right, bottom)))
                if crop.size == 0:
                    continue
                mask = crop.min(axis=2) < 245
                pixels = crop[mask] if mask.any() else crop.reshape(-1, 3)
                rgb = tuple(float(x) for x in pixels.mean(axis=0))
                color = _classify_color_name(rgb)
                cx = left + width / 2
                cy = top + height / 2
                key = (text, round(cx / 24), round(cy / 24))
                current = {
                    "value": text,
                    "color": color,
                    "position": _position_name(cx, cy, rgb_large.width, rgb_large.height),
                    "bbox": {"left": left, "top": top, "width": width, "height": height},
                    "center": {"x": round(cx, 1), "y": round(cy, 1)},
                    "confidence": conf,
                    "variant": variant_name,
                    "rgb": tuple(round(x, 1) for x in rgb),
                }
                prev = candidates.get(key)
                if prev is None or current["confidence"] > prev["confidence"]:
                    candidates[key] = current
        finally:
            if tmp_path and os.path.exists(tmp_path):
                os.unlink(tmp_path)

    return sorted(
        candidates.values(),
        key=lambda item: (-item["confidence"], item["center"]["y"], item["center"]["x"]),
    )


def _format_local_number_extraction(task: str, candidates: List[Dict[str, Any]]) -> Optional[str]:
    """Format OCR/color-grouping output into the observation style used by image_recognition."""
    if len(candidates) < 2:
        return None

    # Only trust short digit groups for direct structured extraction.
    # Long merged strings are useful for debugging, but should fall back to a
    # vision model instead of being treated as clean OCR output.
    filtered = [item for item in candidates if 1 <= len(str(item["value"])) <= 2]
    if len(filtered) < 4:
        return None

    if task == "extract_numbers":
        return "\n".join(
            f"NUMBER: {item['value']}, COLOR: {item['color']}, POSITION: {item['position']}, CONFIDENCE: {item['confidence']:.1f}"
            for item in filtered
        )

    if task == "extract_by_color":
        grouped: Dict[str, List[str]] = {}
        for item in filtered:
            grouped.setdefault(item["color"], []).append(item["value"])
        if len(grouped) < 2:
            return None
        lines: List[str] = []
        for color, values in sorted(grouped.items()):
            lines.append(f"COLOR: {color}")
            lines.append(f"ITEMS: {', '.join(values)}")
            lines.append("")
        return "\n".join(lines).strip()

    return None


def _build_structured_vision_scaffold_prompt(
    task: str,
    candidates: List[Dict[str, Any]],
) -> Optional[str]:
    """Build a verification-first prompt from local OCR/color candidates."""
    if len(candidates) < 2:
        return None

    plausible = [item for item in candidates if 1 <= len(str(item.get("value", ""))) <= 3]
    merged = [item for item in candidates if len(str(item.get("value", ""))) > 3]

    lines = [
        "Use the image itself as the final authority.",
        "You are given OCR/color candidates from a local extractor. They may contain misses, merges, or wrong colors.",
        "Verify them against the image, correct errors, add any missing items, and output only the final structured extraction.",
        "",
        "More plausible OCR/color candidates:",
    ]
    for item in plausible[:16]:
        lines.append(
            "- VALUE: {value}, COLOR: {color}, POSITION: {position}, CONFIDENCE: {confidence:.1f}".format(
                value=item.get("value", ""),
                color=item.get("color", "unknown"),
                position=item.get("position", "unknown"),
                confidence=float(item.get("confidence", 0.0)),
            )
        )
    if merged:
        lines.extend(
            [
                "",
                "Potentially merged OCR chunks (treat as hints, not ground truth):",
            ]
        )
        for item in merged[:12]:
            lines.append(
                "- VALUE: {value}, COLOR: {color}, POSITION: {position}, CONFIDENCE: {confidence:.1f}".format(
                    value=item.get("value", ""),
                    color=item.get("color", "unknown"),
                    position=item.get("position", "unknown"),
                    confidence=float(item.get("confidence", 0.0)),
                )
            )

    if task == "extract_numbers":
        lines.extend(
            [
                "",
                "Return the corrected final extraction in this exact style, one item per line:",
                "NUMBER: [value], COLOR: [color], POSITION: [location]",
                "Do not include explanation or chain-of-thought.",
            ]
        )
    elif task == "extract_by_color":
        lines.extend(
            [
                "",
                "Return the corrected final extraction grouped by color in this exact style:",
                "COLOR: [color]",
                "ITEMS: [comma-separated values]",
                "Do not include explanation or chain-of-thought.",
            ]
        )
    else:
        return None

    return "\n".join(lines)

def image_recognition(
    file_path: str,
    task: str = "describe",
    custom_prompt: str = "",
    output_dir: str = "",
) -> str:
    """
    Image recognition using vision-capable API with automatic fallback.
    Analyzes images by calling the vision API with the image.

    This tool is designed for Text-only LLM planners to extract visual information.
    The vision model acts as a tool backend, returning structured data that enables
    the planner to complete the task.

    Args:
        file_path: Path to the image file
        task: Task type - "describe", "extract_text", "extract_numbers", "extract_by_color",
              "chess", "music_sheet", "geometry", "fractions"
        custom_prompt: Optional custom prompt for specialized analysis. If provided,
                       this overrides the task-based prompt template.
        output_dir: Optional run output directory used to persist debug artifacts.

    Status: ACTIVE
    Dependencies: Vision-capable API (Llama-3.2-90B-Vision, Phi-4 multimodal, etc.)

    Model Priority:
        1. Current planner model, when the endpoint is verified to accept native images
        2. Primary: Llama-3.2-90B-Vision-Instruct (default)
        3. Fallback: Microsoft-Phi-4-multimodal-instruct (if primary fails)
        4. Local OCR + color grouping hints for number/color tasks
           (set VISION_LOCAL_STRUCTURED_FINAL=true to trust them as final output)

    Environment Variables:
        - VISION_MODEL_NAME: Primary vision model (default: Llama-3.2-90B-Vision-Instruct)
        - VISION_FALLBACK_MODEL: Fallback vision model (default: Microsoft-Phi-4-multimodal-instruct)
        - VISION_API_BASE: API endpoint (default: LLM_API_BASE)
        - VISION_API_KEY: API key (default: LLM_API_KEY)
        - VISION_USE_PLANNER_MODEL: Try the current model first if native-vision capable (default: true)
        - VISION_IMAGE_MAX_SIZE: Max image dimension in pixels (default: 1024)
        - VISION_TIMEOUT: API timeout in seconds (default: 300)
        - VISION_MAX_RETRIES: Max retry attempts per model (default: 5)
        - VISION_SKIP_HEALTH_CHECK: Skip health check for faster startup (default: false)
    """
    import time
    import random

    debug_payload: Dict[str, Any] = {
        "file_path": file_path,
        "task": task,
        "custom_prompt": custom_prompt,
        "attempts": [],
    }

    refusal_markers = (
        "can't analyze images",
        "cannot analyze images",
        "unable to analyze images",
        "unable to analyze the image",
        "unable to view images",
        "unable to view any image",
        "unable to view the image",
        "unable to view your image",
        "unable to see the image",
        "unable to see an image",
        "unable to see any image",
        "unable to see your image",
        "unable to see the image you",
        "unable to extract any information because no image was provided",
        "unable to extract any visual information",
        "cannot directly view",
        "cannot view images",
        "cannot view the image",
        "cannot view any image",
        "can't see any image",
        "can't see the image",
        "can't access the image",
        "cannot access the image",
        "don't see any image",
        "do not see any image",
        "don't see an image",
        "do not see an image",
        "don't see the image",
        "do not see the image",
        "don't have access to an image",
        "do not have access to an image",
        "don't have access to the image",
        "do not have access to the image",
        "no image was provided",
        "no image provided",
        "no image attached",
        "no image is attached",
        "none was provided",
        "provide the image",
        "not seeing any image",
        "please provide the image",
        "please upload the image",
        "as an ai text-based model",
        "as a text-only model",
        "provide a description",
        "please describe the image",
        "i'm sorry, but i can't analyze",
    )

    if not os.path.exists(file_path):
        return f"[ERROR] Image not found: {file_path}"

    prompt_family = canonicalize_image_recognition_prompt_family(task, custom_prompt)
    debug_payload["prompt_family"] = prompt_family
    cache_key = _build_file_task_cache_key(file_path, task, prompt_family)
    result_cache = getattr(image_recognition, "_result_cache", {})
    negative_cache = getattr(image_recognition, "_negative_cache", {})
    cached_result = result_cache.get(cache_key)
    if isinstance(cached_result, str) and cached_result:
        print(f"  [INFO] Returning cached image analysis for {os.path.basename(file_path)} [{task}] family={prompt_family}")
        return cached_result
    cached_negative = negative_cache.get(cache_key)
    if isinstance(cached_negative, str) and cached_negative:
        print(f"  [INFO] Returning cached negative image analysis for {os.path.basename(file_path)} [{task}] family={prompt_family}")
        return cached_negative

    def _cache_positive(result_text: str) -> str:
        _store_bounded_cache(result_cache, cache_key, result_text, max_entries=256)
        image_recognition._result_cache = result_cache
        return result_text

    def _cache_negative(result_text: str) -> str:
        _store_bounded_cache(negative_cache, cache_key, result_text, max_entries=256)
        image_recognition._negative_cache = negative_cache
        return result_text

    try:
        local_structured_result: Optional[str] = None
        prefer_local_structured = (
            os.environ.get("VISION_PREFER_LOCAL_STRUCTURED", "true").strip().lower() != "false"
        )
        if prefer_local_structured and task in {"extract_numbers", "extract_by_color"}:
            try:
                local_candidates = _extract_colored_number_candidates(file_path)
                debug_payload["local_structured_candidates"] = local_candidates
                local_result = _format_local_number_extraction(task, local_candidates)
                local_structured_result = local_result
                allow_local_final = (
                    os.environ.get("VISION_LOCAL_STRUCTURED_FINAL", "false").strip().lower() == "true"
                )
                if local_result and allow_local_final:
                    debug_payload["final_source"] = "local_structured_extraction"
                    debug_payload["final_result"] = local_result
                    debug_path = _save_vision_debug_record(output_dir, file_path, task, debug_payload)
                    if debug_path:
                        debug_payload["debug_record_path"] = debug_path
                        print(f"  [INFO] Saved vision debug record: {debug_path}")
                    debug_note = f"Debug Record: {debug_path}\n\n" if debug_path else ""
                    return _cache_positive(
                        f"[Image Analysis - Local Structured Extraction]\n"
                        f"File: {os.path.basename(file_path)}\n"
                        f"Task: {task}\n"
                        f"Method: OCR + color grouping\n\n"
                        f"{debug_note}"
                        f"Result:\n{local_result}"
                    )
                if local_result:
                    debug_payload["local_structured_result_used_as_hint"] = local_result
                scaffold_prompt = _build_structured_vision_scaffold_prompt(task, local_candidates)
                if scaffold_prompt and not custom_prompt:
                    custom_prompt = scaffold_prompt
                    debug_payload["custom_prompt_from_local_candidates"] = scaffold_prompt
            except Exception as local_e:
                debug_payload["local_structured_error"] = str(local_e)
                print(f"  [WARN] Local structured vision extraction failed: {local_e}")

        import base64
        from openai import OpenAI

        # Use environment variables for API configuration (with fallbacks)
        api_base = normalize_api_base(os.environ.get("VISION_API_BASE") or get_api_base())
        api_key = os.environ.get("VISION_API_KEY") or get_api_key()

        # Vision model priority:
        #   current planner model if it natively accepts images,
        #   then the configured dedicated vision model,
        #   then a fallback multimodal model.
        vision_model = os.environ.get("VISION_MODEL_NAME", "Llama-3.2-90B-Vision-Instruct")
        fallback_model = os.environ.get("VISION_FALLBACK_MODEL", "Microsoft-Phi-4-multimodal-instruct")
        planner_model = os.environ.get("LLM_PLANNING_CURRENT_MODEL", "").strip()
        use_planner_model = os.environ.get("VISION_USE_PLANNER_MODEL", "true").strip().lower() != "false"
        debug_payload["planner_model"] = planner_model

        # OPTIMIZATION: Allow skipping health check for faster startup (useful when API is known stable)
        skip_health_check = os.environ.get("VISION_SKIP_HEALTH_CHECK", "false").lower() == "true"

        # Configurable parameters (OPTIMIZED for Cat_C vision tasks)
        # INCREASED: max_image_size to preserve visual details for chess/music/geometry tasks
        max_image_size = int(os.environ.get("VISION_IMAGE_MAX_SIZE", "1024"))
        # INCREASED: timeout for complex visual reasoning tasks
        timeout_seconds = float(os.environ.get("VISION_TIMEOUT", "180.0"))
        # INCREASED: retries with smarter backoff
        max_retries = int(os.environ.get("VISION_MAX_RETRIES", "4"))

        if not api_key:
            return "[ERROR] VISION_API_KEY or LLM_API_KEY not set. Vision API unavailable."

        # Encode and compress image to base64
        import io
        from PIL import Image

        with Image.open(file_path) as img:
            original_size = img.size
            original_mode = img.mode

            # Force convert to RGB
            if img.mode != "RGB":
                img = img.convert("RGB")

            # OPTIMIZED: Adaptive image sizing based on task type
            # Chess/geometry tasks need higher resolution for piece/shape detection
            if task in ["chess", "geometry", "music_sheet", "extract_numbers"]:
                effective_max_size = min(max_image_size, 1280)  # Higher res for detail tasks
            else:
                effective_max_size = min(max_image_size, 1024)  # Standard for other tasks

            if max(img.size) > effective_max_size:
                factor = effective_max_size / max(img.size)
                new_size = (int(img.width * factor), int(img.height * factor))
                img = img.resize(new_size, Image.Resampling.LANCZOS)
                print(f"  [INFO] Resized image from {original_size} to {new_size} for task: {task}")

            # OPTIMIZED: Adaptive quality based on task type
            # Higher quality for tasks requiring fine detail detection
            if task in ["chess", "music_sheet", "extract_text"]:
                jpeg_quality = 90  # Higher quality for text/symbol recognition
            else:
                jpeg_quality = 80  # Balanced quality for other tasks

            img_byte_arr = io.BytesIO()
            img.save(img_byte_arr, format="JPEG", quality=jpeg_quality)
            img_b64 = base64.b64encode(img_byte_arr.getvalue()).decode("utf-8")

            print(f"  [INFO] Image encoded: {len(img_b64)} bytes (original mode: {original_mode})")

        # The OpenAI spec requires data URI
        image_url = f"data:image/jpeg;base64,{img_b64}"

        # Enhanced system instruction for vision model (FIX: added guidance prompts)
        system_instruction = """You are a vision analysis tool being called by a planning LLM.
Your role is to extract ACCURATE, COMPLETE, and STRUCTURED information from images.

CRITICAL RULES:
1. Extract ALL relevant information systematically (do not skip or summarize)
2. Preserve exact values (numbers, text, colors, positions)
3. Use clear, structured format that the planner can parse
4. If uncertain, state what you see without guessing
5. Focus on extracting DATA, not solving the task

The planning LLM will use your output to perform calculations and reasoning."""

        # Build task-specific prompt (custom_prompt takes priority)
        if custom_prompt:
            prompt = f"{system_instruction}\n\n{custom_prompt}"
        elif task == "describe":
            prompt = f"""{system_instruction}

TASK: Comprehensive Image Description

Describe everything you see in this image in detail. Include:
- Objects and their positions (left, right, center, top, bottom)
- Colors with precision (e.g., "red", "dark green", "light blue")
- All visible text and numbers (transcribe exactly)
- Any patterns, relationships, or notable features

Be systematic and thorough. The planning LLM will use this description to understand the image."""

        elif task == "extract_text" or "text" in task.lower():
            prompt = f"""{system_instruction}

TASK: Text Extraction

Extract and list ALL text visible in this image.
- Include numbers, labels, and any written content
- Preserve the spatial layout if relevant
- Note the font style or color if it appears significant
- Output format: List each text element on a new line

Do NOT interpret or solve problems - just extract the text."""

        elif task == "extract_numbers" or "number" in task.lower():
            prompt = f"""{system_instruction}

TASK: Number Extraction with Color Information

Extract ALL numbers from this image with their colors.

For EACH number visible:
1. State the number value (exact)
2. State its EXACT color (be precise: "red", "green", "blue", etc.)
3. Note its position if relevant

Output format (one per line):
NUMBER: [value], COLOR: [color], POSITION: [location]

Example:
NUMBER: 42, COLOR: red, POSITION: top-left
NUMBER: 13, COLOR: green, POSITION: center

This color-number association is CRITICAL for the planner's calculations.
Extract ALL numbers - do not skip any."""

        elif task == "extract_by_color":
            prompt = f"""{system_instruction}

TASK: Organize Content by Color

Analyze this image and organize ALL content BY COLOR.

For each distinct color present:
COLOR: [exact color name]
ITEMS: [list all numbers/text in that color, comma-separated]

Example:
COLOR: red
ITEMS: 24, 32, 43, 51

COLOR: green
ITEMS: 13, 22, 31, 41

Be precise with color identification - this is critical for calculations.
List ALL items for each color - do not omit any."""
        elif task == "chess" or "chess" in task.lower():
            prompt = f"""{system_instruction}

TASK: Chess Position Analysis

Analyze this chess position carefully.

1. BOARD ORIENTATION: Identify which side is which (white pieces typically start at bottom)
2. PIECE POSITIONS: For EACH piece on the board, state:
   - Piece type (K=King, Q=Queen, R=Rook, B=Bishop, N=Knight, P=Pawn)
   - Color (White or Black)
   - Square (using algebraic notation: a1-h8)
3. WHOSE TURN: Determine whose turn it is if indicated
4. KEY FEATURES: Note any checks, threats, or tactical patterns
5. POSITION STRING: If possible, provide a FEN-like board string and state any uncertainty.
6. BEST MOVE: If the task asks for a move, return candidate moves in standard SAN notation
   (for example Rd5, Qxg2#, Nf3). Do not use coordinates unless SAN is impossible.

Format your piece list as:
White: Ka1, Qd1, Ra1, Rh1, ...
Black: Ke8, Qd8, Ra8, Rh8, ...
FEN: [piece placement / active color / castling / en-passant / halfmove / fullmove if known]
Candidate SAN moves: [ranked list with confidence]

Be precise with square positions - files are a-h (left to right from White's view),
ranks are 1-8 (bottom to top from White's view).
List ALL pieces - do not omit any. If orientation or a piece is ambiguous, explicitly mark it uncertain."""

        elif task == "music_sheet" or "music" in task.lower():
            prompt = f"""{system_instruction}

TASK: Music Sheet Analysis

Analyze this music sheet notation.

1. CLEF: Identify the clef (treble/bass clef)
2. KEY SIGNATURE: Note any sharps or flats
3. TIME SIGNATURE: Note the time signature if visible
4. NOTES: For each note, identify:
   - Pitch (letter name: A, B, C, D, E, F, G)
   - Octave position (relative to the staff)
   - Duration (whole, half, quarter, eighth, etc.)
   - Position on staff (line or space number)
5. DERIVED TEXT/NUMBER: If the task asks to decode notes into words, ages, or numbers,
   provide both the raw note sequence and the decoded answer candidate.

For TREBLE CLEF (lines bottom to top): E, G, B, D, F
For BASS CLEF (lines bottom to top): G, B, D, F, A

List the notes in order from left to right. Extract ALL notes and preserve accidentals/rests."""

        elif task == "geometry" or "polygon" in task.lower():
            prompt = f"""{system_instruction}

TASK: Geometric Figure Analysis

Analyze this geometric figure.

1. SHAPE: Identify the type of shape/polygon
2. VERTICES: Count the number of vertices/corners
3. LABELED MEASUREMENTS: For each labeled side or angle:
   - State the measurement value (exact)
   - Identify which side/angle it refers to
4. COLOR CODING: Note if different colors indicate different elements
5. DIMENSIONS: List ALL dimensions needed for calculations

List all visible measurements and their locations systematically.
Do not calculate area - just extract the measurements."""

        elif task == "fractions" or "math" in task.lower():
            prompt = f"""{system_instruction}

TASK: Math Worksheet Analysis

Analyze this math worksheet or problem set.

1. PROBLEMS: List each problem visible (transcribe exactly)
2. For each fraction problem:
   - State the fractions involved
   - Identify the operation (add, subtract, multiply, divide)
   - Note any student answers written
3. ANSWERS: If answers are shown, list them
4. GRADING: If this is graded work, note correct/incorrect marks
5. ORDER: Preserve the original left-to-right, top-to-bottom order exactly.

Format: Problem N: [fraction operation], Answer: [answer if shown]
Do not solve - just extract what is visible."""

        else:
            prompt = f"""{system_instruction}

TASK: {task}

Analyze this image and provide detailed information relevant to the task.
Be specific about:
- Positions (use terms like top-left, center, bottom-right)
- Colors (be precise with color names)
- Numbers and text (transcribe exactly)
- Any patterns or relationships visible

Extract information systematically - do not skip details."""

        # OPTIMIZED: Adaptive retry strategy
        base_retry_delay = 3.0  # Reduced initial delay
        last_error = None
        total_attempts = 0

        # Determine max_tokens based on task (OPTIMIZED for task complexity)
        task_token_config = {
            "chess": 6000,  # Chess needs detailed piece positions
            "music_sheet": 5000,  # Music notation can be verbose
            "extract_by_color": 4096,  # Color-organized lists
            "describe": 4096,  # Full description
            "geometry": 3500,  # Measurements and calculations
            "fractions": 3000,  # Math problems
            "extract_numbers": 2500,  # Number lists
            "extract_text": 3000,  # Text extraction
        }
        max_tokens = task_token_config.get(task, 3072)

        models_to_try: List[str] = []

        def add_model_candidate(candidate: Optional[str]) -> None:
            if not candidate:
                return
            normalized_candidate = str(candidate).strip()
            if not normalized_candidate:
                return
            if normalized_candidate not in models_to_try:
                models_to_try.append(normalized_candidate)

        if use_planner_model and planner_model and supports_native_vision_input(planner_model):
            add_model_candidate(planner_model)
        add_model_candidate(vision_model)
        add_model_candidate(fallback_model)
        debug_payload["model_priority"] = models_to_try

        # OPTIMIZATION: Conditional health check with caching
        # Skip health check when: (1) explicitly disabled, (2) cached result available
        _vision_health_cache = getattr(image_recognition, '_health_cache', {})
        _cache_expiry = getattr(image_recognition, '_cache_expiry', 0)
        current_time = time.time()

        healthy_models = []

        if skip_health_check:
            print(f"  [INFO] Health check skipped (VISION_SKIP_HEALTH_CHECK=true)")
            healthy_models = models_to_try
        elif current_time < _cache_expiry and _vision_health_cache:
            # Use cached health check results (valid for 5 minutes)
            healthy_models = [m for m in models_to_try if _vision_health_cache.get(m, False)]
            if healthy_models:
                print(f"  [INFO] Using cached health check results: {healthy_models}")
            else:
                healthy_models = models_to_try  # Cache invalid, try all
        else:
            # Perform fresh health check with SHORTER timeout
            for test_model in models_to_try:
                try:
                    # Create tiny diagnostic image. A plain "OK?" health check can
                    # pass even when the endpoint ignores image payloads, so ask
                    # for a visual property that is not stated in the prompt.
                    import io as io_test
                    from PIL import Image as Image_test
                    test_img = Image_test.new('RGB', (32, 32), color='blue')
                    for x in range(16):
                        for y in range(32):
                            test_img.putpixel((x, y), (255, 0, 0))
                    test_byte_arr = io_test.BytesIO()
                    test_img.save(test_byte_arr, format="JPEG", quality=30)
                    test_b64 = base64.b64encode(test_byte_arr.getvalue()).decode("utf-8")
                    test_url = f"data:image/jpeg;base64,{test_b64}"

                    # Quick health check (3s timeout - reduced from 5s)
                    test_client = OpenAI(
                        base_url=api_base,
                        api_key=api_key,
                        default_headers={"User-Agent": "curl/7.68.0"},
                        timeout=3.0,
                    )
                    test_response = test_client.chat.completions.create(
                        model=test_model,
                        messages=[{"role": "user", "content": [
                            {"type": "text", "text": "What color is the left half of the attached image? Reply with one color word."},
                            {"type": "image_url", "image_url": {"url": test_url}}
                        ]}],
                        **_chat_completion_sampling_kwargs(test_model),
                        **_chat_completion_token_kwargs(test_model, 8),
                    )
                    test_result = (test_response.choices[0].message.content or "").strip().lower()
                    if (
                        not test_result
                        or any(marker in test_result for marker in refusal_markers)
                        or "red" not in test_result
                    ):
                        raise RuntimeError(f"vision health check did not confirm image access: {test_result[:80]}")
                    healthy_models.append(test_model)
                    _vision_health_cache[test_model] = True
                    print(f"  [INFO] Health check passed for vision model: {test_model}")
                except Exception as e:
                    error_msg = str(e)[:100]
                    _vision_health_cache[test_model] = False
                    # IMPROVED: Only log as warning if it's a blocking error, not timeout
                    if "blocked" in error_msg.lower():
                        print(f"  [WARN] Vision model {test_model} blocked by API, skipping")
                    elif "timeout" in error_msg.lower() or "504" in error_msg or "502" in error_msg:
                        print(f"  [WARN] Vision model {test_model} timed out during health check, will retry on actual call")
                        healthy_models.append(test_model)  # Still add - might work on real call
                    else:
                        print(f"  [WARN] Health check failed for {test_model}: {error_msg}")
                    continue

            # Update cache expiry (5 minutes)
            image_recognition._health_cache = _vision_health_cache
            image_recognition._cache_expiry = current_time + 300

        # If all health checks failed, proceed with original list anyway (fallback)
        if healthy_models:
            models_to_try = healthy_models
        else:
            print(f"  [WARN] All health checks failed, trying original models anyway")
            models_to_try = list(debug_payload.get("model_priority") or models_to_try)

        for model_idx, current_model in enumerate(models_to_try):
            is_fallback = model_idx > 0
            if is_fallback:
                print(f"  [INFO] Trying fallback vision model: {current_model}")

            for attempt in range(max_retries):
                try:
                    # FIX: Use configurable timeout (increased from 180 to 300)
                    client = OpenAI(
                        base_url=api_base,
                        api_key=api_key,
                        default_headers={"User-Agent": "curl/7.68.0"},
                        timeout=timeout_seconds,
                    )

                    # FIX: Use configurable model name instead of hardcoded
                    response = client.chat.completions.create(
                        model=current_model,
                        messages=[{
                            "role": "user",
                            "content": [
                                {"type": "text", "text": prompt},
                                {
                                    "type": "image_url",
                                    "image_url": {"url": image_url}
                                }
                            ]
                        }],
                        **_chat_completion_sampling_kwargs(current_model),
                        **_chat_completion_token_kwargs(current_model, max_tokens),
                    )

                    result = (response.choices[0].message.content or "").strip()
                    debug_payload["attempts"].append({
                        "model": current_model,
                        "attempt": attempt + 1,
                        "status": "ok",
                        "response": result,
                    })

                    if not result:
                        last_error = RuntimeError("vision_refusal: empty_response")
                        print(f"  [WARN] Vision API ({current_model}) returned empty content, trying fallback if available")
                        break

                    if any(marker in result.lower() for marker in refusal_markers):
                        refusal_snippet = result.splitlines()[0][:200]
                        last_error = RuntimeError(f"vision_refusal: {refusal_snippet}")
                        print(f"  [WARN] Vision API ({current_model}) returned no usable visual analysis, trying fallback if available")
                        break

                    # Success - return formatted result
                    model_note = f" (fallback)" if is_fallback else ""
                    debug_payload["final_source"] = current_model + model_note
                    debug_payload["final_result"] = result
                    debug_path = _save_vision_debug_record(output_dir, file_path, task, debug_payload)
                    if debug_path:
                        debug_payload["debug_record_path"] = debug_path
                        print(f"  [INFO] Saved vision debug record: {debug_path}")
                    debug_note = f"Debug Record: {debug_path}\n\n" if debug_path else ""
                    return _cache_positive(
                        f"[Image Analysis]\n"
                        f"File: {os.path.basename(file_path)}\n"
                        f"Task: {task}\n"
                        f"Model: {current_model}{model_note}\n\n"
                        f"{debug_note}"
                        f"Result:\n{result}"
                    )

                except Exception as e:
                    last_error = e
                    error_str = str(e).lower()
                    total_attempts += 1
                    debug_payload["attempts"].append({
                        "model": current_model,
                        "attempt": attempt + 1,
                        "status": "error",
                        "error": str(e),
                    })

                    # ENHANCED: More granular error classification
                    gateway_errors = ["504", "502", "503", "gateway"]
                    timeout_errors = ["timeout", "timed out", "read timeout"]
                    connection_errors = ["connection", "connect", "network"]
                    rate_limit_errors = ["rate limit", "too many requests", "429"]
                    capacity_errors = ["overloaded", "capacity", "busy"]
                    blocked_errors = ["blocked", "forbidden", "403"]

                    # Determine error type and retry strategy
                    is_gateway = any(err in error_str for err in gateway_errors)
                    is_timeout = any(err in error_str for err in timeout_errors)
                    is_connection = any(err in error_str for err in connection_errors)
                    is_rate_limit = any(err in error_str for err in rate_limit_errors)
                    is_capacity = any(err in error_str for err in capacity_errors)
                    is_blocked = any(err in error_str for err in blocked_errors)

                    # Non-retryable: blocked errors -> skip to next model immediately
                    if is_blocked:
                        print(f"  [WARN] Vision API ({current_model}) request blocked, skipping to fallback")
                        break

                    should_retry = is_gateway or is_timeout or is_connection or is_rate_limit or is_capacity

                    if should_retry and attempt < max_retries - 1:
                        # ADAPTIVE backoff: shorter for timeouts, longer for rate limits
                        if is_rate_limit:
                            exponential_delay = base_retry_delay * (3 ** attempt)  # Longer for rate limits
                        elif is_gateway or is_timeout:
                            exponential_delay = base_retry_delay * (1.5 ** attempt)  # Shorter for gateway issues
                        else:
                            exponential_delay = base_retry_delay * (2 ** attempt)

                        jitter = random.uniform(0, min(exponential_delay * 0.25, 8.0))
                        wait_time = min(exponential_delay + jitter, 60.0)  # Cap at 60 seconds

                        error_type = "gateway" if is_gateway else "timeout" if is_timeout else "connection" if is_connection else "rate_limit" if is_rate_limit else "capacity"
                        print(f"  [WARN] Vision API ({current_model}) {error_type} error (attempt {attempt + 1}/{max_retries})")
                        print(f"  [WARN] Retrying in {wait_time:.1f}s...")
                        time.sleep(wait_time)
                        continue

                    # Non-retryable error or max retries reached - try fallback
                    if attempt == max_retries - 1:
                        print(f"  [ERROR] Vision API ({current_model}) failed after {max_retries} attempts")
                    break  # Break inner loop, try next model

        # All Vision API models failed - try enhanced local analysis as final fallback
        print(f"  [INFO] All Vision APIs failed after {total_attempts} total attempts, trying local analysis...")

        if local_structured_result and task in {"extract_numbers", "extract_by_color"}:
            debug_payload["final_source"] = "local_structured_extraction_after_vision_failure"
            debug_payload["final_result"] = local_structured_result
            debug_path = _save_vision_debug_record(output_dir, file_path, task, debug_payload)
            if debug_path:
                debug_payload["debug_record_path"] = debug_path
                print(f"  [INFO] Saved vision debug record: {debug_path}")
            debug_note = f"Debug Record: {debug_path}\n\n" if debug_path else ""
            return _cache_positive(
                f"[Image Analysis - Local Structured Extraction]\n"
                f"File: {os.path.basename(file_path)}\n"
                f"Task: {task}\n"
                f"Method: OCR + color grouping after vision fallback\n\n"
                f"{debug_note}"
                f"Result:\n{local_structured_result}"
            )

        try:
            from PIL import Image as PILImage
            import pytesseract
            pytesseract.pytesseract.tesseract_cmd = _find_tesseract_binary()

            ocr_img = PILImage.open(file_path)
            width, height = ocr_img.size
            mode = ocr_img.mode

            result_parts = [
                f"[Image Analysis - Local Fallback]",
                f"File: {os.path.basename(file_path)}",
                f"Task: {task}",
                f"Note: Vision API unavailable, using local analysis (limited capability)",
                f"",
                f"Image Info: {width}x{height}, mode={mode}",
            ]

            # ENHANCED: Try OCR with multiple configurations
            ocr_text = ""
            try:
                # Standard OCR
                ocr_text = pytesseract.image_to_string(ocr_img)

                # If standard fails, try with different PSM modes for specific tasks
                if not ocr_text.strip() and task in ["extract_numbers", "extract_text"]:
                    # PSM 6: Assume a single uniform block of text
                    ocr_text = pytesseract.image_to_string(ocr_img, config='--psm 6')

                if not ocr_text.strip() and task == "extract_numbers":
                    # PSM 7: Treat the image as a single text line (good for numbers)
                    ocr_text = pytesseract.image_to_string(ocr_img, config='--psm 7 -c tessedit_char_whitelist=0123456789.,-')

            except Exception as ocr_e:
                print(f"  [WARN] OCR extraction failed: {str(ocr_e)[:50]}")

            if ocr_text.strip():
                result_parts.append(f"\nExtracted Text/Numbers:\n{ocr_text.strip()}")
                print(f"  [INFO] Local OCR extracted {len(ocr_text.strip())} characters")
            else:
                result_parts.append("\nNo text detected by OCR.")

            # ENHANCED: Basic color analysis for color-related tasks
            if task in ["extract_by_color", "extract_numbers", "describe"]:
                try:
                    from collections import Counter

                    # Convert to RGB if needed
                    if ocr_img.mode != "RGB":
                        ocr_img = ocr_img.convert("RGB")

                    # Sample pixels and count dominant colors
                    pixels = list(ocr_img.getdata())
                    # Quantize colors to reduce noise (group similar colors)
                    def quantize_color(rgb, levels=8):
                        step = 256 // levels
                        return tuple((c // step) * step for c in rgb)

                    quantized = [quantize_color(p) for p in pixels[:10000]]  # Sample first 10k pixels
                    color_counts = Counter(quantized)
                    top_colors = color_counts.most_common(5)

                    def rgb_to_name(rgb):
                        r, g, b = rgb
                        if r > 200 and g < 100 and b < 100:
                            return "red"
                        elif r < 100 and g > 200 and b < 100:
                            return "green"
                        elif r < 100 and g < 100 and b > 200:
                            return "blue"
                        elif r > 200 and g > 200 and b < 100:
                            return "yellow"
                        elif r > 200 and g > 200 and b > 200:
                            return "white"
                        elif r < 50 and g < 50 and b < 50:
                            return "black"
                        elif r > 200 and g < 150 and b > 200:
                            return "purple"
                        elif r > 200 and g > 100 and b < 100:
                            return "orange"
                        else:
                            return f"rgb({r},{g},{b})"

                    color_info = [f"{rgb_to_name(c)}: {count} pixels" for c, count in top_colors]
                    result_parts.append(f"\nDominant Colors (approximate):\n" + "\n".join(color_info))
                except Exception as color_e:
                    print(f"  [WARN] Color analysis failed: {str(color_e)[:50]}")

            result_parts.append("\n[WARNING] Local fallback cannot detect: complex shapes, chess pieces, music notes, or precise spatial relationships. Results may be incomplete for this task type.")
            local_fallback_result = "\n".join(result_parts)
            debug_payload["final_source"] = "local_fallback"
            debug_payload["final_result"] = local_fallback_result
            debug_path = _save_vision_debug_record(output_dir, file_path, task, debug_payload)
            if debug_path:
                debug_payload["debug_record_path"] = debug_path
                print(f"  [INFO] Saved vision debug record: {debug_path}")
                local_fallback_result = local_fallback_result.replace(
                    "[Image Analysis - Local Fallback]",
                    f"[Image Analysis - Local Fallback]\nDebug Record: {debug_path}",
                    1,
                )
            return _cache_positive(local_fallback_result)

        except ImportError as imp_err:
            print(f"  [WARN] pytesseract not installed: {imp_err}, OCR fallback unavailable")
        except Exception as ocr_error:
            print(f"  [WARN] Local analysis failed: {str(ocr_error)[:100]}")

        # All methods failed - return detailed error message
        if last_error and str(last_error).startswith("vision_refusal:"):
            refusal_reason = str(last_error).split("vision_refusal:", 1)[1].strip() or "unusable_response"
            error_result = (
                f"[ERROR] Vision model returned no usable visual analysis.\n"
                f"File: {os.path.basename(file_path)}\n"
                f"Task: {task}\n"
                f"Reason: {refusal_reason}"
            )
            debug_payload["final_source"] = "vision_refusal"
            debug_payload["final_result"] = error_result
            debug_path = _save_vision_debug_record(output_dir, file_path, task, debug_payload)
            if debug_path:
                debug_payload["debug_record_path"] = debug_path
                print(f"  [INFO] Saved vision debug record: {debug_path}")
                error_result = f"{error_result}\nDebug Record: {debug_path}"
            return _cache_negative(error_result)

        error_type = type(last_error).__name__ if last_error else "Unknown"
        models_tried = ", ".join(models_to_try)
        error_result = f"[ERROR] Image recognition failed after trying all models ({models_tried}) and local OCR.\nError type: {error_type}\nLast error: {str(last_error)}"
        debug_payload["final_source"] = "error"
        debug_payload["final_result"] = error_result
        debug_path = _save_vision_debug_record(output_dir, file_path, task, debug_payload)
        if debug_path:
            debug_payload["debug_record_path"] = debug_path
            print(f"  [INFO] Saved vision debug record: {debug_path}")
            error_result = f"{error_result}\nDebug Record: {debug_path}"
        return error_result

    except ImportError as e:
        return f"[ERROR] Missing dependency: {e}. Install with: pip install openai"
    except Exception as e:
        return f"[ERROR] Image recognition failed: {str(e)}"



def video_analysis(video_path: str, task: str = "general") -> str:
    """Mock video analysis - would use video processing libraries."""
    if not os.path.exists(video_path):
        return f"[MOCK] Video not found: {video_path}"
    size = os.path.getsize(video_path)
    return f"[MOCK] Video analysis of {os.path.basename(video_path)} ({size} bytes). Task: {task}. (Would require opencv/moviepy)"



def reasoning(problem: str) -> str:
    """
    Perform API-backed reasoning on a logic-only subproblem.

    This tool is intentionally separate from the planner turn: the planner can
    call it after collecting observations, and the helper returns a concise
    calculation/check rather than a fake acknowledgement.
    """
    problem = str(problem or "").strip()
    if not problem:
        return "[ERROR] reasoning requires a non-empty problem."

    try:
        from openai import OpenAI

        api_base = normalize_api_base(os.environ.get("REASONING_API_BASE") or get_api_base())
        api_key = os.environ.get("REASONING_API_KEY") or get_api_key()
        model_name = (
            os.environ.get("REASONING_MODEL_NAME")
            or os.environ.get("LLM_PLANNING_CURRENT_MODEL")
            or os.environ.get("LLM_API_MODEL")
            or "Llama-3.3-70B-Instruct"
        ).strip()
        timeout_seconds = float(os.environ.get("REASONING_TIMEOUT", "120"))
        max_tokens = int(os.environ.get("REASONING_MAX_TOKENS", "1536"))

        if not api_key:
            return "[ERROR] REASONING_API_KEY or LLM_API_KEY not set. Reasoning API unavailable."
        if not model_name:
            return "[ERROR] No reasoning model configured."

        prompt = f"""You are a deterministic reasoning helper inside a tool-use benchmark.
Use only the facts provided in the problem statement and previously observed tool outputs.
Do not browse, invent missing evidence, or claim to inspect files/images directly.

Return a compact analysis with:
1. Given facts
2. Calculation or inference
3. Conclusion

Problem:
{problem}
"""
        client = OpenAI(
            base_url=api_base,
            api_key=api_key,
            default_headers={"User-Agent": "curl/7.68.0"},
            timeout=timeout_seconds,
        )
        response = client.chat.completions.create(
            model=model_name,
            messages=[
                {
                    "role": "system",
                    "content": "You are a concise reasoning tool. Produce plain text, no JSON required.",
                },
                {"role": "user", "content": prompt},
            ],
            **_chat_completion_sampling_kwargs(model_name),
            **_chat_completion_token_kwargs(model_name, max_tokens),
        )
        result = (response.choices[0].message.content or "").strip()
        if not result:
            return f"[ERROR] Reasoning model {model_name} returned empty output."
        return f"[Reasoning Analysis]\nModel: {model_name}\n\n{result}"
    except ImportError as e:
        return f"[ERROR] Missing dependency for reasoning tool: {e}. Install with: pip install openai"
    except Exception as e:
        return f"[ERROR] Reasoning tool failed: {type(e).__name__}: {e}"

def code_interpreter(code: str, language: str = "python", action: str = "execute") -> str:
    """Mock code interpreter - for non-Python languages."""
    if language.lower() == "python":
        return python_executor(code)
    return f"[MOCK] Code interpreter for {language}. Action: {action}. Code:\n{code[:200]}..."



# =============================================================================
# Active Tool 9: download_file
# =============================================================================

def download_file(url: str, save_path: Optional[str] = None, output_dir: str = DEFAULT_ARTIFACT_DIR) -> str:
    """
    Download a file from a URL using requests.
    Useful for obtaining PDF/Excel files before reading them.

    Status: IMPLEMENTED
    Dependencies: requests
    """
    try:
        import requests
    except ImportError:
        return "[ERROR] requests not installed. Run: pip install requests"

    try:
        url = _sanitize_url_input(url)
        downloaded = _download_url_to_path(
            url=url,
            save_path=save_path,
            output_dir=output_dir,
            timeout=30,
        )
        return (
            f"[SUCCESS] Downloaded {downloaded['requested_url']}\n"
            f"Resolved URL: {downloaded['final_url']}\n"
            f"Saved to: {downloaded['save_path']}\n"
            f"Content-Type: {downloaded['content_type'] or 'unknown'}\n"
            f"Size: {downloaded['size']} bytes\n"
            "(You can now use file_reader, pdf_reader, or excel_reader on this path)"
        )

    except Exception as e:
        return f"[ERROR] Failed to download {url}: {str(e)}"


# =============================================================================
# Tool Implementation Map
# =============================================================================




# =============================================================================
# Active Tool 12: submit_final_answer
# =============================================================================

def normalize_submitted_answer(answer: Any) -> Tuple[Any, Optional[str]]:
    """Best-effort unwrap for JSON-wrapped or container-form final answers."""

    placeholder_markers = (
        "the final concise answer string",
        "final concise answer string",
        "final answer string",
        "placeholder answer",
        "placeholder final answer",
        "replace with the final answer",
        "your final answer here",
        "insert final answer",
    )

    def _looks_like_placeholder_answer(value: str) -> bool:
        normalized = re.sub(r"\s+", " ", str(value or "").strip().lower())
        if not normalized:
            return False
        if normalized in {"answer", "final answer", "result", "output", "value", "response"}:
            return True
        return any(marker in normalized for marker in placeholder_markers)

    def _strip_code_fence(text: str) -> str:
        stripped = text.strip()
        if stripped.startswith("```") and stripped.endswith("```"):
            lines = stripped.splitlines()
            if len(lines) >= 3:
                return "\n".join(lines[1:-1]).strip()
        return stripped

    def _extract(value: Any) -> Tuple[Any, Optional[str]]:
        if isinstance(value, str):
            stripped = _strip_code_fence(value)
            if not stripped:
                return stripped, None
            if _looks_like_placeholder_answer(stripped):
                return "", "placeholder_answer_rejected"
            looks_structured = (
                (stripped.startswith("{") and stripped.endswith("}"))
                or (stripped.startswith("[") and stripped.endswith("]"))
                or (stripped.startswith('"') and stripped.endswith('"'))
            )
            if looks_structured:
                try:
                    parsed = json.loads(stripped)
                except Exception:
                    return stripped, None
                extracted, note = _extract(parsed)
                return extracted, note or "json_wrapped_answer_unwrapped"
            return stripped, None

        if isinstance(value, list):
            cleaned_items = []
            for item in value:
                extracted, _ = _extract(item)
                if isinstance(extracted, str) and extracted.strip():
                    cleaned_items.append(extracted.strip())
            if cleaned_items:
                return ", ".join(cleaned_items), "list_answer_joined"
            return value, None

        if isinstance(value, dict):
            for key in ("answer", "final_answer", "response", "output", "result"):
                if key in value:
                    extracted, note = _extract(value[key])
                    if isinstance(extracted, str) and extracted.strip():
                        return extracted, note or f"dict_{key}_unwrapped"

            non_empty = {k: v for k, v in value.items() if v not in (None, "", [], {})}
            if len(non_empty) == 1:
                only_key, only_value = next(iter(non_empty.items()))
                extracted, note = _extract(only_value)
                if isinstance(extracted, str) and extracted.strip():
                    return extracted, note or f"single_key_dict_{only_key}_unwrapped"
            return value, None

        return value, None

    return _extract(answer)


def submit_final_answer(answer: str, answer_type: str = "string") -> str:
    """
    Submit the final answer to the task.
    This tool marks the task as completed.
    YOU MUST USE THIS TOOL TO FINISH THE TASK.

    Args:
        answer: The final answer string.
        answer_type: The type of answer (e.g., "string", "number", "date"). Default: string.
    """
    normalized_answer, _ = normalize_submitted_answer(answer)
    if isinstance(normalized_answer, (dict, list)):
        try:
            normalized_answer = json.dumps(normalized_answer, ensure_ascii=False)
        except Exception:
            normalized_answer = str(normalized_answer)
    return f"[FINAL ANSWER SUBMITTED] {normalized_answer} (Type: {answer_type})"


# =============================================================================
# Tool Implementation Map
# =============================================================================

TOOL_IMPLEMENTATIONS = {
    # Active Tools
    "web_search": web_search,
    "excel_reader": excel_reader,
    "audio_transcription": audio_transcription,
    "file_reader": file_reader,
    "python_executor": python_executor,
    "pptx_reader": pptx_reader,
    "pdf_reader": pdf_reader,
    "zip_extractor": zip_extractor,
    "download_file": download_file,
    "web_browser": web_browser,
    "calculator": calculator,
    "submit_final_answer": submit_final_answer,

    # Optional / Mock Tools
    "image_recognition": image_recognition,
    "video_analysis": video_analysis,
    "reasoning": reasoning, 
    "code_interpreter": code_interpreter,
}


def execute_tool(
    tool_name: str,
    tool_args: Dict[str, Any],
    attachments_dir: str = "",
    output_dir: str = ""
) -> str:
    """
    Execute a tool by name with given arguments.
    Automatically resolves file paths relative to attachments_dir.

    Args:
        tool_name: Name of the tool to execute
        tool_args: Arguments for the tool
        attachments_dir: Directory containing input files/attachments
        output_dir: Directory for tool outputs (plots, extracted files, etc.)
    """
    if tool_name not in TOOL_IMPLEMENTATIONS:
        return f"[ERROR] Unknown tool: {tool_name}. Available: {list(TOOL_IMPLEMENTATIONS.keys())}"

    tool_args = dict(tool_args or {})

    # Map legacy parameter names to normalized file_path
    if "image_path" in tool_args and "file_path" not in tool_args:
        tool_args["file_path"] = tool_args.pop("image_path")
    if "audio_path" in tool_args and "file_path" not in tool_args:
        tool_args["file_path"] = tool_args.pop("audio_path")
    if tool_name == "pdf_reader" and "url" in tool_args and "file_path" not in tool_args:
        tool_args["file_path"] = tool_args.pop("url")
    if tool_name == "web_browser" and "file_path" in tool_args and "url" not in tool_args:
        tool_args["url"] = tool_args.pop("file_path")

    # Lightweight argument alias recovery for common single-argument tools.
    if tool_name == "reasoning" and "problem" not in tool_args:
        for alias in ("query", "question", "task", "prompt", "text", "input"):
            if alias in tool_args and tool_args[alias]:
                tool_args["problem"] = tool_args.pop(alias)
                break
        if "problem" not in tool_args and len(tool_args) == 1:
            only_value = next(iter(tool_args.values()))
            if isinstance(only_value, str) and only_value.strip():
                tool_args = {"problem": only_value}

    if tool_name == "python_executor" and "code" not in tool_args:
        for alias in ("python", "python_code", "script", "program", "snippet", "expression"):
            if alias in tool_args and tool_args[alias]:
                candidate = tool_args.pop(alias)
                if alias == "expression":
                    candidate = f"result = {candidate}\nprint(result)"
                tool_args["code"] = candidate
                break
        if "code" not in tool_args and len(tool_args) == 1:
            only_value = next(iter(tool_args.values()))
            if isinstance(only_value, str) and only_value.strip():
                tool_args = {"code": only_value}

    resolution_notes: List[str] = []

    # Normalize URL arguments early so all URL tools see the same cleaned target.
    if tool_name in {"web_browser", "download_file", "pdf_reader"} and "url" in tool_args:
        tool_args["url"] = _sanitize_url_input(tool_args["url"])
    if tool_name in {"web_browser", "download_file", "pdf_reader"} and "file_path" in tool_args:
        value = str(tool_args.get("file_path") or "")
        if _is_http_url(_sanitize_url_input(value)):
            tool_args["file_path"] = _sanitize_url_input(value)

    # Resolve file paths
    path_keys = ["file_path", "video_path"]
    for key in path_keys:
        if key in tool_args:
            resolved, note = _resolve_file_argument(tool_args[key], attachments_dir)
            tool_args[key] = resolved
            if note:
                resolution_notes.append(note)

    # Auto-route common attachment types when the model picks file_reader too broadly.
    if tool_name == "file_reader" and tool_args.get("file_path"):
        ext = Path(str(tool_args["file_path"])).suffix.lower()
        routed_tool = None
        routed_args = {"file_path": tool_args["file_path"]}
        if ext == ".pdf":
            routed_tool = "pdf_reader"
        elif ext in {".csv", ".xlsx", ".xls"}:
            routed_tool = "excel_reader"
            if tool_args.get("query"):
                routed_args["query"] = tool_args["query"]
        elif ext == ".pptx":
            routed_tool = "pptx_reader"
        elif ext == ".zip":
            routed_tool = "zip_extractor"
            if tool_args.get("extract_to"):
                routed_args["extract_to"] = tool_args["extract_to"]
        elif ext in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif", ".tif", ".tiff"}:
            routed_tool = "image_recognition"
            hint = (
                tool_args.get("task")
                or tool_args.get("query")
                or tool_args.get("question")
                or tool_args.get("prompt")
                or tool_args.get("custom_prompt")
                or ""
            )
            hint_text = str(hint or "")
            hint_lower = hint_text.lower()
            if tool_args.get("task"):
                routed_args["task"] = str(tool_args["task"])
            elif any(marker in hint_lower for marker in ("text", "ocr", "read", "transcribe", "extract")):
                routed_args["task"] = "extract_text"
            else:
                routed_args["task"] = "describe"
            if tool_args.get("custom_prompt"):
                routed_args["custom_prompt"] = str(tool_args["custom_prompt"])
            elif hint_text and not tool_args.get("task"):
                routed_args["custom_prompt"] = hint_text

        if routed_tool:
            tool_name = routed_tool
            tool_args = routed_args

    # Tools that support output_dir parameter
    tools_with_output_dir = {"python_executor", "zip_extractor", "download_file", "image_recognition"}

    try:
        func = TOOL_IMPLEMENTATIONS[tool_name]

        # Pass runtime context to tools that can use it.
        if tool_name == "python_executor":
            sanitized_args = dict(tool_args)
            sanitized_args.pop("output_dir", None)
            sanitized_args.pop("attachments_dir", None)
            if output_dir:
                sanitized_args["output_dir"] = output_dir
            if attachments_dir:
                runtime_attachments_dir, runtime_input_files = _format_input_files_for_python(attachments_dir)
                sanitized_args["attachments_dir"] = runtime_attachments_dir
                sanitized_args.setdefault("input_files", runtime_input_files)
            result = func(**sanitized_args)
        elif tool_name in tools_with_output_dir and output_dir:
            sanitized_args = dict(tool_args)
            sanitized_args.pop("output_dir", None)
            result = func(**sanitized_args, output_dir=output_dir)
        else:
            result = func(**tool_args)

        if resolution_notes:
            return "\n".join(f"[INFO] {note}" for note in resolution_notes) + "\n\n" + str(result)
        return result

    except TypeError as e:
        return f"[ERROR] Invalid arguments for {tool_name}: {str(e)}"
    except Exception as e:
        return f"[ERROR] Tool execution failed: {str(e)}"


# =============================================================================
# Test / Demo
# =============================================================================

if __name__ == "__main__":
    print("Tool Execution Layer for GAIA Mini")
    print("=" * 50)

    # Count active vs mock tools
    active_tools = []
    mock_tools = []
    for name, func in TOOL_IMPLEMENTATIONS.items():
        doc = func.__doc__.split('\n')[0] if func.__doc__ else ""
        if "[MOCK]" in doc or "Mock" in doc:
            mock_tools.append(name)
        else:
            active_tools.append(name)

    print(f"Active Tools ({len(active_tools)}): {active_tools}")
    print(f"Mock Tools ({len(mock_tools)}): {mock_tools}")
    print("\nAll available tools:")
    for name, func in TOOL_IMPLEMENTATIONS.items():
        doc = func.__doc__.split('\n')[0] if func.__doc__ else "No description"
        status = "MOCK" if "[MOCK]" in doc or "Mock" in doc else "ACTIVE"
        print(f"  [{status}] {name}: {doc[:60]}...")
