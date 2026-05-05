#!/usr/bin/env python3
"""
GAIA-Specific Utilities

This module provides:
1. File path resolution and attachment reading for GAIA benchmark
2. Data format conversion from raw GAIA to unified evaluation format

The code supports both:
- categorized local mirrors such as data/GAIA/cat_B_document/attachments/<uuid>.xlsx
- snapshot-style layouts such as .../2023/validation/<uuid>.xlsx
"""
from __future__ import annotations
import os
import re
from pathlib import Path
from typing import Dict, Any, Optional, List


_PROJECT_ROOT = Path(__file__).resolve().parents[2]
_LEGACY_SNAPSHOT_DIR = _PROJECT_ROOT / "data" / "datasets--gaia-benchmark--GAIA" / "snapshots" / "682dd723ee1e1697e00360edccf2366dc8418dd9" / "2023"
_GAIA_CATEGORY_BY_SUFFIX = {
    "A": "cat_A_text",
    "B": "cat_B_document",
    "C": "cat_C_vision",
    "D": "cat_D_audio",
}

# Reading tool IDs that should return extracted content
GAIA_READING_TOOLS = {
    "pdf_reader", "excel_reader", "file_reader", "pptx_reader",
    "pdb_analyzer", "zip_extractor", "image_recognition",
    "audio_transcription", "video_analysis"
}


def _get_gaia_roots() -> List[Path]:
    """Return candidate GAIA dataset roots in search order."""
    roots: List[Path] = []

    data_root = os.getenv("GAIA_DATA_ROOT")
    if data_root:
        roots.append(Path(data_root).expanduser())

    env_root = os.getenv("GAIA_BASE_DIR")
    if env_root:
        roots.append(Path(env_root).expanduser())

    roots.extend([
        _PROJECT_ROOT.parent / "data" / "Augmented",
        _PROJECT_ROOT.parent / "data.hg" / "Augmented",
        _PROJECT_ROOT.parent / "data" / "GAIA",
        _PROJECT_ROOT.parent / "data.hg" / "GAIA",
        _PROJECT_ROOT / "data" / "Augmented",
        _PROJECT_ROOT / "data.hg" / "Augmented",
        _PROJECT_ROOT / "data" / "GAIA",
        _PROJECT_ROOT / "data.hg" / "GAIA",
        _LEGACY_SNAPSHOT_DIR,
    ])

    deduped: List[Path] = []
    seen = set()
    for root in roots:
        key = str(root)
        if key in seen:
            continue
        seen.add(key)
        deduped.append(root)
    return deduped


def _resolve_against_gaia_roots(path_obj: Path) -> Optional[str]:
    """Resolve paths containing a GAIA root marker against configured GAIA roots."""
    path_parts = path_obj.parts
    if "GAIA" not in path_parts:
        return None
    gaia_idx = path_parts.index("GAIA")
    gaia_relative = Path(*path_parts[gaia_idx + 1:])
    for root in _get_gaia_roots():
        rooted = root / gaia_relative
        if rooted.exists():
            return str(rooted.resolve())
    return None


def _infer_gaia_category_dir(record: Optional[Dict[str, Any]] = None) -> Optional[str]:
    """Infer the categorized GAIA subdirectory from record metadata."""
    if not record:
        return None
    subset = str(record.get("meta", {}).get("subset", "") or "")
    match = re.search(r"_([ABCD])$", subset)
    if not match:
        return None
    return _GAIA_CATEGORY_BY_SUFFIX.get(match.group(1))


def resolve_attachment_path(
    file_path: Optional[str] = None,
    file_name: Optional[str] = None,
    record: Optional[Dict[str, Any]] = None,
) -> Optional[str]:
    """
    Resolve a GAIA attachment path robustly.

    Priority:
    1. Existing provided file_path
    2. Project-root-relative interpretation of a relative file_path
    3. file_name lookup through known GAIA roots
    """
    candidate_name = file_name or ""

    if file_path:
        path_obj = Path(file_path).expanduser()
        rooted_path = _resolve_against_gaia_roots(path_obj)
        if rooted_path:
            return rooted_path

        if path_obj.exists():
            return str(path_obj.resolve())

        if not path_obj.is_absolute():
            rel = (_PROJECT_ROOT / path_obj).resolve()
            if rel.exists():
                return str(rel)

            rooted_path = _resolve_against_gaia_roots(path_obj)
            if rooted_path:
                return rooted_path

        if path_obj.name and not candidate_name:
            candidate_name = path_obj.name

    if candidate_name:
        return resolve_gaia_file_path(candidate_name, record=record)
    return None


def get_attachment_display_path(
    attachment: Dict[str, Any],
    record: Optional[Dict[str, Any]] = None,
) -> str:
    """Return a current-machine-safe attachment path for prompts/logging."""
    resolved = resolve_attachment_path(
        attachment.get("file_path"),
        attachment.get("file_name"),
        record=record,
    )
    if resolved:
        resolved_path = Path(resolved)
        try:
            return str(resolved_path.relative_to(_PROJECT_ROOT))
        except ValueError:
            return str(resolved_path)
    return attachment.get("file_path") or attachment.get("file_name") or ""


def resolve_gaia_file_path(
    file_name: str,
    split: str = "validation",
    record: Optional[Dict[str, Any]] = None
) -> Optional[str]:
    """
    Resolve the absolute file path for a GAIA attachment.

    GAIA attachments are UUID-named files stored in split-specific subfolders:
    - Base: /path/to/GAIA/snapshots/.../2023/
    - Path: Base + split (validation/test) + UUID_filename

    Args:
        file_name: UUID-named filename (e.g., "076c8171-9b3b-49b9-a477-244d2a532826.xlsx")
        split: Dataset split ("validation" or "test")
        record: Optional record dict to extract split from meta

    Returns:
        Absolute file path if found, None otherwise
    """
    if not file_name:
        return None

    # Try to get split from record metadata if available
    if record:
        split = record.get("meta", {}).get("split", split)

    category_dir = _infer_gaia_category_dir(record)

    for root in _get_gaia_roots():
        # First try categorized local mirrors: <root>/cat_B_document/attachments/<file>
        if category_dir:
            categorized_path = root / category_dir / "attachments" / file_name
            if categorized_path.exists():
                return str(categorized_path.resolve())
        else:
            for candidate_category_dir in _GAIA_CATEGORY_BY_SUFFIX.values():
                categorized_path = root / candidate_category_dir / "attachments" / file_name
                if categorized_path.exists():
                    return str(categorized_path.resolve())

        # Then try snapshot-style mirrors: <root>/<split>/<file>
        direct_path = root / split / file_name
        if direct_path.exists():
            return str(direct_path.resolve())

        for fallback_split in ["validation", "test"]:
            fallback_path = root / fallback_split / file_name
            if fallback_path.exists():
                return str(fallback_path.resolve())

    return None


def read_attachment_content(
    attachment: Dict[str, Any],
    max_chars: int = 10000,
    record: Optional[Dict[str, Any]] = None
) -> Optional[str]:
    """
    Read attachment file content for GAIA dataset with split-aware path resolution.

    Supports multi-format reading:
    - xlsx/xls/csv: Read as table text (pandas)
    - txt/json/jsonld/py: Read as plain text
    - pdf: Read text content (pypdf or pdfplumber)
    - docx: Read Word documents (python-docx)
    - pptx: Read PowerPoint presentations (python-pptx)
    - pdb: Read Protein Data Bank files
    - zip: List archive contents

    Args:
        attachment: Dict with file_name, file_path, file_type
        max_chars: Maximum characters to return (truncate if longer)
        record: Optional record dict for split-aware path resolution

    Returns:
        File content as string, or None if cannot read
    """
    file_path = attachment.get("file_path")
    file_type = attachment.get("file_type", "").lower()
    file_name = attachment.get("file_name", "")

    # Split-aware and machine-safe path resolution for GAIA
    resolved_path = resolve_attachment_path(file_path, file_name, record=record)
    if resolved_path:
        file_path = resolved_path
    else:
        return None

    try:
        # Excel files
        if file_type in ["xlsx", "xls", "csv"]:
            try:
                import pandas as pd
                if file_type == "csv":
                    df = pd.read_csv(file_path)
                else:
                    df = pd.read_excel(file_path)
                content = f"[Excel/CSV File: {file_name}]\n"
                content += f"Shape: {df.shape[0]} rows x {df.shape[1]} columns\n"
                content += f"Columns: {list(df.columns)}\n\n"
                content += df.to_string(max_rows=100, max_cols=20)
                return content[:max_chars]
            except Exception as e:
                return f"[Error reading Excel file: {e}]"

        # Text files
        elif file_type in ["txt", "json", "jsonld", "py", "md", "xml", "html"]:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            return f"[Text File: {file_name}]\n{content[:max_chars]}"

        # PDF files
        elif file_type == "pdf":
            try:
                # Try pypdf first
                try:
                    from pypdf import PdfReader
                    reader = PdfReader(file_path)
                    text = ""
                    for page in reader.pages[:10]:  # First 10 pages
                        text += page.extract_text() + "\n"
                    return f"[PDF File: {file_name}]\n{text[:max_chars]}"
                except ImportError:
                    pass

                # Try pdfplumber
                try:
                    import pdfplumber
                    with pdfplumber.open(file_path) as pdf:
                        text = ""
                        for page in pdf.pages[:10]:
                            text += (page.extract_text() or "") + "\n"
                    return f"[PDF File: {file_name}]\n{text[:max_chars]}"
                except ImportError:
                    return f"[PDF File: {file_name}] - Cannot read (install pypdf or pdfplumber)"
            except Exception as e:
                return f"[Error reading PDF: {e}]"

        # PDB (Protein Data Bank) files
        elif file_type == "pdb":
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            return f"[PDB File: {file_name}]\n{content[:max_chars]}"

        # DOCX files
        elif file_type == "docx":
            try:
                from docx import Document
                doc = Document(file_path)
                text = "\n".join([para.text for para in doc.paragraphs])
                return f"[Word Document: {file_name}]\n{text[:max_chars]}"
            except ImportError:
                return f"[Word Document: {file_name}] - Cannot read (install python-docx)"
            except Exception as e:
                return f"[Error reading DOCX: {e}]"

        # PPTX files
        elif file_type == "pptx":
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return f"[PowerPoint: {file_name}]\n{text[:max_chars]}"
            except ImportError:
                return f"[PowerPoint: {file_name}] - Cannot read (install python-pptx)"
            except Exception as e:
                return f"[Error reading PPTX: {e}]"

        # ZIP files - list contents
        elif file_type == "zip":
            try:
                import zipfile
                with zipfile.ZipFile(file_path, 'r') as zf:
                    names = zf.namelist()
                    content = f"[ZIP Archive: {file_name}]\nContents ({len(names)} files):\n"
                    for name in names[:50]:
                        content += f"  - {name}\n"
                    return content[:max_chars]
            except Exception as e:
                return f"[Error reading ZIP: {e}]"

        # Image files
        elif file_type in ["png", "jpg", "jpeg", "gif", "bmp"]:
            return f"[Image File: {file_name}] - Image content cannot be read as text. Describe the visual information needed."

        # Audio files
        elif file_type in ["mp3", "wav", "m4a", "ogg"]:
            return f"[Audio File: {file_name}] - Audio content cannot be transcribed in this context."

        # Video files
        elif file_type in ["mp4", "mov", "avi", "mkv"]:
            return f"[Video File: {file_name}] - Video content cannot be analyzed in this context."

        else:
            # Try reading as text
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                return f"[File: {file_name}]\n{content[:max_chars]}"
            except:
                return f"[File: {file_name}] - Cannot read content"

    except Exception as e:
        return f"[Error reading {file_name}: {e}]"


def execute_reading_tool(
    tool_id: str,
    arguments: Dict[str, Any],
    record: Optional[Dict[str, Any]] = None
) -> Optional[str]:
    """
    Execute a GAIA reading tool and return the extracted content.

    When the model outputs a tool_call for a reading tool (pdf_reader, excel_reader,
    file_reader, etc.), this function returns the content extracted from the
    attachment file resolved in the GAIA dataset.

    Args:
        tool_id: The tool ID being called (e.g., "pdf_reader", "excel_reader")
        arguments: Tool arguments (should contain file_path or file_name)
        record: The GAIA record for split-aware path resolution

    Returns:
        Extracted content string if successful, None otherwise
    """
    if tool_id not in GAIA_READING_TOOLS:
        return None

    # Extract file path/name from arguments
    file_path = None
    file_name = None

    if isinstance(arguments, dict):
        file_path = arguments.get("file_path") or arguments.get("path")
        file_name = arguments.get("file_name") or arguments.get("name")
    elif isinstance(arguments, list):
        for arg in arguments:
            if isinstance(arg, dict):
                if arg.get("name") in ["file_path", "path"]:
                    file_path = arg.get("value")
                elif arg.get("name") in ["file_name", "name"]:
                    file_name = arg.get("value")

    # If no file info in arguments, try to get from record attachments
    if not file_path and not file_name and record:
        attachments = record.get("query", {}).get("attachments", [])
        if attachments:
            att = attachments[0]  # Use first attachment
            file_name = att.get("file_name")
            file_path = att.get("file_path")

    # Resolve path and read content
    resolved_path = resolve_attachment_path(file_path, file_name, record=record)
    if resolved_path:
        file_path = resolved_path

    if not file_path:
        return None

    # Determine file type from extension
    file_type = ""
    if file_path:
        ext = os.path.splitext(file_path)[1].lower().lstrip(".")
        file_type = ext

    # Create attachment dict and read
    attachment = {
        "file_path": file_path,
        "file_type": file_type,
        "file_name": file_name or os.path.basename(file_path)
    }

    return read_attachment_content(attachment, record=record)


# =============================================================================
# GAIA Data Format Conversion (Solution to 0-score problem)
# =============================================================================

def parse_gaia_tools_string(tools_str: str) -> List[Dict[str, Any]]:
    """
    Parse GAIA tools string format to structured tool list.

    GAIA format example:
        "submit_final_answer(answer: str) -> str
         web_search(query: str) -> str
         calculator(expression: str) -> str"

    Returns:
        List of tool dicts with tool_id, description, arguments_schema
    """
    tools = []

    if not tools_str or not isinstance(tools_str, str):
        return tools

    # Split by tool definitions (function signatures)
    # Pattern: function_name(args) -> return_type followed by description
    lines = tools_str.strip().split('\n')

    current_tool_id = None
    current_desc = []

    for line in lines:
        line = line.strip()
        if not line:
            continue

        # Check if line starts with a function signature
        func_match = re.match(r'^(\w+)\((.*?)\)\s*(?:->.*)?$', line)

        if func_match:
            # Save previous tool if exists
            if current_tool_id:
                tools.append({
                    "tool_id": current_tool_id,
                    "description": '\n'.join(current_desc).strip(),
                    "arguments_schema": {}  # GAIA doesn't provide detailed schema
                })

            # Start new tool
            current_tool_id = func_match.group(1)
            current_desc = []
        else:
            # Description line
            current_desc.append(line)

    # Add last tool
    if current_tool_id:
        tools.append({
            "tool_id": current_tool_id,
            "description": '\n'.join(current_desc).strip(),
            "arguments_schema": {}
        })

    return tools


def convert_gaia_record_to_unified(gaia_record: Dict[str, Any]) -> Dict[str, Any]:
    """
    Convert raw GAIA format to unified evaluation format.

    This fixes the 0-score issue by:
    1. Adding proper meta.dataset field for evaluator recognition
    2. Creating gold.final_answer structure from raw final_answer string
    3. Ensuring proper query structure with user_query field

    Args:
        gaia_record: Raw GAIA record with task_id, question, final_answer, etc.

    Returns:
        Unified format record with meta, query, tool_environment, gold fields
    """
    # Extract basic fields
    task_id = gaia_record.get("task_id", "")
    question = gaia_record.get("question", "")
    final_answer = gaia_record.get("final_answer", "")
    level = gaia_record.get("level", "")
    tools_str = gaia_record.get("tools", "")

    # Parse tools
    tools = parse_gaia_tools_string(tools_str)

    # Handle file_name field (GAIA-specific attachment format)
    attachments = []
    if "file_name" in gaia_record and gaia_record["file_name"]:
        file_name = gaia_record["file_name"]
        # Resolve file path
        file_path = resolve_gaia_file_path(file_name)
        if file_path:
            # Determine file type from extension
            file_type = os.path.splitext(file_name)[1].lstrip('.').lower()
            attachments.append({
                "file_name": file_name,
                "file_path": file_path,
                "file_type": file_type
            })

    # Build unified format
    unified = {
        "meta": {
            "dataset": "gaia",
            "id": task_id,
            "level": level,
            "has_arguments": False,  # GAIA only provides final answers, no gold tool calls
            "split": "validation"  # Default to validation, can be overridden
        },
        "query": {
            "user_query": question,
            "attachments": attachments
        },
        "tool_environment": {
            "tools": tools
        },
        "gold": {
            "final_answer": {
                "answer_type": "string",
                "answer": final_answer,
                "aliases": []
            },
            # GAIA doesn't provide gold plan_dag/tool_calls in the dataset
            # These will be empty for answer-mode evaluation
            "plan_dag": {
                "nodes": [],
                "edges": []
            },
            "tool_calls": []
        }
    }

    # Preserve original fields for reference
    unified["_original"] = {
        "task_id": task_id,
        "original_tools": gaia_record.get("original_tools", ""),
        "annotator_metadata": gaia_record.get("annotator_metadata", {})
    }

    return unified


def is_gaia_format(record: Dict[str, Any]) -> bool:
    """
    Detect if a record is in raw GAIA format (needs conversion).

    Raw GAIA format has:
    - task_id field (not "meta" or "query")
    - question field (not "user_query")
    - final_answer as string (not dict)

    Unified format has:
    - meta.dataset field
    - query.user_query field
    - gold.final_answer.answer field
    """
    # Check for raw GAIA indicators
    has_task_id = "task_id" in record
    has_question = "question" in record
    has_flat_answer = isinstance(record.get("final_answer"), str)

    # Check for unified format indicators
    has_meta = "meta" in record and isinstance(record.get("meta"), dict)
    has_query = "query" in record and isinstance(record.get("query"), dict)

    # If has raw GAIA fields and lacks unified fields, it's raw GAIA
    if has_task_id and has_question and has_flat_answer:
        if not has_meta or not has_query:
            return True

    return False
